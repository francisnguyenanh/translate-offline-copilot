# -*- coding: utf-8 -*-
"""
Ứng dụng Flask quản lý trích xuất và nạp bản dịch cho file Excel, PowerPoint và Word
"""

import os
import json
import zipfile
import uuid
import shutil
import io
import re
import csv
import requests as _requests
from datetime import datetime, timedelta
from urllib.parse import quote
from flask import Flask, render_template, request, send_file, send_from_directory, jsonify, session, redirect, url_for, Response, stream_with_context
from werkzeug.utils import secure_filename
from copy import deepcopy, copy
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document
from functools import wraps
from lxml import etree as _etree

# Khởi tạo ứng dụng Flask
app = Flask(__name__, static_folder='templates/static', static_url_path='/static')
app.config['MAX_CONTENT_LENGTH'] = 900 * 1024 * 1024  # Giới hạn 900MB

# DEBUG: enable debug-level logging for OCR coordinate diagnostics
import logging
app.logger.setLevel(logging.DEBUG)
logging.basicConfig(level=logging.DEBUG)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['SECRET_KEY'] = os.urandom(24)  # Secret key cho session
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(hours=5)  # Session timeout 5h

# Các định dạng file được phép
ALLOWED_EXTENSIONS = {'xlsx', 'pptx', 'docx'}

# Đọc password từ file
PASSWORD_FILE = 'password.txt'


# File lưu Prompt Templates (Tab 1/3)
TEMPLATES_FILE = 'prompt_templates.json'
# File lưu Prompt Template cho Tab 2 (Dịch ảnh)
IMG_OCR_PROMPT_FILE = 'img_ocr_prompt_template.json'
def load_img_ocr_prompt_template():
    """Đọc prompt template cho Tab 2 (Dịch ảnh OCR)"""
    try:
        with open(IMG_OCR_PROMPT_FILE, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if isinstance(data, list):
            return data
        return [data]
    except Exception:
        # Fallback mặc định nếu chưa có file
        return [{
            "id": "img-ocr-default",
            "name": "Prompt dịch ảnh OCR (default)",
            "content": "Translate the following numbered text items from {sourceLang} to {targetLang}.\nReturn ONLY the same numbered list with translated text. Do not add any other text.\n\n{listText}"
        }]

def save_img_ocr_prompt_template(new_templates):
    """Ghi prompt template cho Tab 2 (Dịch ảnh OCR)"""
    with open(IMG_OCR_PROMPT_FILE, 'w', encoding='utf-8') as f:
        json.dump(new_templates, f, ensure_ascii=False, indent=2)

GLOSSARY_DIR = 'glossaries'   # thư mục lưu các file CSV chuyên ngành
os.makedirs(GLOSSARY_DIR, exist_ok=True)

# ==================== GOOGLE SHEET HELPERS ====================

def parse_google_sheet_url(url: str):
    """
    Trả về (spreadsheet_id, gid) từ link Google Sheet bất kỳ.
    Hỗ trợ: .../spreadsheets/d/{ID}/edit#gid={GID} hoặc dạng khác.
    """
    match = re.search(r'/spreadsheets/d/([a-zA-Z0-9_-]+)', url)
    if not match:
        raise ValueError('URL không hợp lệ. Vui lòng dùng link Google Sheet.')
    spreadsheet_id = match.group(1)
    gid_match = re.search(r'[#&?]gid=(\d+)', url)
    gid = gid_match.group(1) if gid_match else '0'
    return spreadsheet_id, gid


def build_sheet_export_url(spreadsheet_id: str) -> str:
    """Tạo URL export toàn bộ file sang .xlsx (không cần auth nếu file public)."""
    return (f'https://docs.google.com/spreadsheets/d/{spreadsheet_id}'
            f'/export?format=xlsx&id={spreadsheet_id}')

# ==================== HELPER: PROMPT TEMPLATES ====================
def get_default_templates():
    """Trả về danh sách template mặc định"""
    return [
        {
            "id": "formal",
            "name": "Dịch chính xác (Formal)",
            "content": "Hãy dịch các giá trị (values) trong file JSON này sang {TARGET_LANG}.\n\nPhong cách: Chính xác, chuyên nghiệp, dùng trong tài liệu kinh doanh.\n\nQuy tắc bắt buộc:\n1. Giữ nguyên 100% các keys\n2. CHỈ dịch nội dung bên trong values\n3. KHÔNG dịch từ/cụm từ đã là ngôn ngữ đích\n4. KHÔNG dịch mã kỹ thuật, placeholder, tên biến\n5. KHÔNG dịch số, ngày tháng, ký hiệu đặc biệt\n6. Giữ nguyên format JSON chuẩn\n\n⚠️ QUY TẮc về dấu ngoặc kép: CHỈ dùng \" (U+0022). KHÔNG dùng \u201c \u201d \u201e \u201f \u00ab \u00bb\nTrích dẫn: dùng \u300c \u300dhoặc 'đơn'\n\nOutput: Trả về ĐÚNG cấu trúc JSON, KHÔNG thêm giải thích."
        },
        {
            "id": "casual",
            "name": "Dịch tự nhiên (Casual)",
            "content": "Hãy dịch các giá trị (values) trong file JSON này sang {TARGET_LANG}.\n\nPhong cách: Tự nhiên, thân thiện, dễ đọc - phù hợp cho giao diện người dùng.\n\nQuy tắc bắt buộc:\n1. Giữ nguyên 100% các keys\n2. CHỈ dịch nội dung bên trong values\n3. KHÔNG dịch từ/cụm từ đã là ngôn ngữ đích\n4. KHÔNG dịch mã kỹ thuật, placeholder, tên biến\n5. KHÔNG dịch số, ngày tháng, ký hiệu đặc biệt\n6. Giữ nguyên format JSON chuẩn\n\n⚠️ QUY TẮc về dấu ngoặc kép: CHỈ dùng \" (U+0022). KHÔNG dùng \u201c \u201d \u201e \u201f \u00ab \u00bb\n\nOutput: Trả về ĐÚNG cấu trúc JSON, KHÔNG thêm giải thích."
        },
        {
            "id": "technical",
            "name": "Dịch kỹ thuật (Technical)",
            "content": "Hãy dịch các giá trị (values) trong file JSON này sang {TARGET_LANG}.\n\nPhong cách: Kỹ thuật, chính xác cao, giữ nguyên thuật ngữ IT.\n\nQuy tắc bắt buộc:\n1. Giữ nguyên 100% các keys\n2. CHỈ dịch nội dung bên trong values\n3. KHÔNG dịch từ/cụm từ đã là ngôn ngữ đích\n4. KHÔNG dịch placeholder ({0}, %s, $n...), tên biến\n5. KHÔNG dịch số, ngày tháng, ký hiệu đặc biệt\n6. Giữ nguyên thuật ngữ IT tiếng Anh nếu không có từ tương đương chính xác\n7. Giữ nguyên format JSON chuẩn\n\n⚠️ QUY TẮc về dấu ngoặc kép: CHỈ dùng \" (U+0022). KHÔNG dùng \u201c \u201d \u201e \u201f \u00ab \u00bb\n\nOutput: Trả về ĐÚNG cấu trúc JSON, KHÔNG thêm giải thích."
        }
    ]

def load_templates(lang='default'):
    """Đọc prompt templates cho một ngôn ngữ cụ thể (fallback về default)"""
    try:
        with open(TEMPLATES_FILE, 'r', encoding='utf-8') as f:
            data = json.load(f)
        # Tương thích ngược: nếu data là array thì đó là format cũ
        if isinstance(data, list):
            return data
        return data.get(lang) or data.get('default') or get_default_templates()
    except Exception:
        return get_default_templates()


def apply_glossary(extracted_data: dict, glossary_ids: list) -> dict:
    """
    Thay thế các cụm từ trong extracted_data theo các file glossary được chọn.

    Logic:
    - Với mỗi cặp (src, dst) trong glossary: tìm src (case-insensitive, exact match)
      trong từng value của extracted_data, thay bằng dst.
    - "Exact match" = cụm từ đứng độc lập, không phải substring nằm giữa ký tự chữ.
      Dùng regex word-boundary \\b kết hợp re.IGNORECASE.
    - Ưu tiên thay thế cụm dài trước (tránh thay một phần của cụm dài hơn).
    - Nếu glossary_ids rỗng hoặc không có file nào, trả về nguyên extracted_data.

    CSV format: cột A = ngôn ngữ đích (dst), cột B = ngôn ngữ gốc (src)
    """
    if not glossary_ids:
        return extracted_data

    # Thu thập tất cả cặp (src, dst) từ các file glossary được chọn
    pairs = []
    for gid in glossary_ids:
        filepath = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
        if not os.path.exists(filepath):
            continue
        with open(filepath, 'r', encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for row in reader:
                if len(row) >= 2:
                    dst = row[0].strip()   # cột A: ngôn ngữ đích
                    src = row[1].strip()   # cột B: ngôn ngữ gốc
                    if src and dst:
                        pairs.append((src, dst))

    if not pairs:
        return extracted_data

    # Sắp xếp: cụm dài trước để tránh thay nhầm phần của cụm dài
    pairs.sort(key=lambda x: len(x[0]), reverse=True)

    # Build regex patterns (một lần, tái sử dụng)
    # Dùng \b nếu src bắt đầu/kết thúc bằng ký tự word; fallback lookaround nếu không
    def make_pattern(src: str):
        escaped = re.escape(src)
        prefix = r'\b' if re.match(r'\w', src[0])  else r'(?<![^\s])'
        suffix = r'\b' if re.match(r'\w', src[-1]) else r'(?![^\s])'
        return re.compile(prefix + escaped + suffix, re.IGNORECASE)

    compiled = [(make_pattern(src), dst) for src, dst in pairs]

    # Áp dụng thay thế lên từng value
    result = {}
    for key, value in extracted_data.items():
        if not isinstance(value, str):
            result[key] = value
            continue
        for pattern, dst in compiled:
            value = pattern.sub(dst, value)
        result[key] = value

    return result


def build_dedup_data(extracted_data, chunk_size=400):
    """
    Gộp các keys có cùng value để giảm số lượng cần dịch.
    Returns: (dedup_files, mapping, stats)
      - dedup_files: list of {name, content} – các chunk dedup (giống format files thường)
      - mapping: {dedup_key: [orig_key1, orig_key2, ...]}
      - stats: {total, unique, saved, percent_saved}
    """
    # Group keys by value (giữ order)
    value_to_keys = {}
    for key, value in extracted_data.items():
        if value not in value_to_keys:
            value_to_keys[value] = []
        value_to_keys[value].append(key)

    # Build dedup dict và mapping
    dedup_data = {}
    mapping = {}  # dedup_key → [original_keys]
    for idx, (value, keys) in enumerate(value_to_keys.items(), 1):
        dk = f'dedup_{idx}'
        dedup_data[dk] = value
        mapping[dk] = keys

    total = len(extracted_data)
    unique = len(dedup_data)
    saved = total - unique
    percent = round(saved * 100 / total) if total > 0 else 0
    stats = {'total': total, 'unique': unique, 'saved': saved, 'percent_saved': percent}

    # Chia thành các chunk
    items = list(dedup_data.items())
    num_chunks = max(1, (unique + chunk_size - 1) // chunk_size)
    dedup_files = []
    for i in range(num_chunks):
        chunk = dict(items[i * chunk_size:(i + 1) * chunk_size])
        dedup_files.append({
            'name': f'dedup_part{i+1:02d}_of_{num_chunks:02d}.json',
            'content': json.dumps(chunk, ensure_ascii=False, indent=2)
        })

    return dedup_files, mapping, stats


def expand_dedup_data(json_data, session_folder):
    """
    Mở rộng dedup JSON (dedup_N → value) thành keys gốc dựa trên mapping đã lưu.
    Nếu không tìm thấy mapping file thì trả về nguyên.
    """
    mapping_path = os.path.join(session_folder, 'dedup_mapping.json')
    if not os.path.exists(mapping_path):
        return json_data
    try:
        with open(mapping_path, 'r', encoding='utf-8') as f:
            mapping = json.load(f)
        expanded = {}
        for key, value in json_data.items():
            if key.startswith('dedup_') and key in mapping:
                for orig_key in mapping[key]:
                    expanded[orig_key] = value
            else:
                expanded[key] = value  # key thường, giữ nguyên
        return expanded
    except Exception:
        return json_data


def stream_extract(filepath, original_filename, glossary_ids, session_folder, color_filter=None, proofread_mode=False):
    """
    Generator cho SSE progress events khi trích xuất file.
    Yields chuỗi SSE format: data: {json}\n\n
    """
    def _evt(step, pct, **kwargs):
        payload = {'step': step, 'pct': pct, **kwargs}
        return f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"

    try:
        original_ext = original_filename.rsplit('.', 1)[-1].lower() if '.' in original_filename else 'xlsx'
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        yield _evt('reading', 10, message='Đang đọc file...')

        # Bước 1: Trích xuất text
        if original_ext == 'xlsx':
            workbook = load_workbook(filepath)
            extracted_data = {}
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        if isinstance(cell.value, str) and not cell.value.startswith('='):
                            if color_filter is None or (_get_font_rgb_xlsx(cell) or '000000') in color_filter:
                                extracted_data[f"{sheet_name}!{cell.coordinate}"] = cell.value
            extracted_data.update(extract_xlsx_shapes(filepath))
            workbook.close()
        elif original_ext == 'pptx':
            extracted_data = extract_text_from_pptx(filepath, color_filter)
        elif original_ext == 'docx':
            extracted_data = extract_text_from_docx(filepath, color_filter)
        else:
            yield _evt('error', 0, error=f'Không hỗ trợ định dạng .{original_ext}')
            return

        if proofread_mode:
            extracted_data = _filter_proofread_extract_data(extracted_data)

        yield _evt('chunking', 40, message='Đang áp dụng glossary...')

        if glossary_ids:
            extracted_data = apply_glossary(extracted_data, glossary_ids)

        yield _evt('chunking', 60, message='Đang tạo file JSON...')

        # Bước 2: Chia thành chunks
        CHUNK_SIZE = 300
        data_items = list(extracted_data.items())
        total_items = len(data_items)
        num_files = max(1, (total_items + CHUNK_SIZE - 1) // CHUNK_SIZE)

        base_filename = os.path.splitext(original_filename)[0] or f'file_{timestamp}'
        safe_base = f'extracted_{timestamp}'
        folder_name = f'{base_filename}_json_to_translate'
        temp_dir = os.path.join(session_folder, f'{safe_base}_temp_{timestamp}')
        os.makedirs(temp_dir, exist_ok=True)

        json_files = []
        json_display_names = []
        for i in range(num_files):
            chunk_data = dict(data_items[i * CHUNK_SIZE:(i + 1) * CHUNK_SIZE])
            disp_name = f'{base_filename}_part{i+1:02d}_of_{num_files:02d}.json'
            safe_name = f'{safe_base}_part{i+1:02d}.json'
            jpath = os.path.join(temp_dir, safe_name)
            with open(jpath, 'w', encoding='utf-8') as jf:
                json.dump(chunk_data, jf, ensure_ascii=False, indent=2)
            json_files.append(jpath)
            json_display_names.append(disp_name)

        files_data = []
        for idx, jpath in enumerate(json_files):
            with open(jpath, 'r', encoding='utf-8') as f:
                files_data.append({'name': json_display_names[idx], 'content': f.read()})

        yield _evt('writing', 75, message='Đang tạo ZIP và dedup...')

        # Tạo ZIP
        zip_display = f'{base_filename}_json_to_translate.zip'
        safe_zip = f'{safe_base}_json_{timestamp}.zip'
        zip_path = os.path.join(session_folder, safe_zip)
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_STORED) as zf:
            for idx, jpath in enumerate(json_files):
                zf.write(jpath, os.path.join(folder_name, json_display_names[idx]))

        # Ghi trạng thái ra file (không thể ghi session từ trong generator)
        extract_state = {
            'path': zip_path,
            'display_name': zip_display,
            'input_path': filepath,
            'json_files': json_files,
            'temp_dir': temp_dir,
        }
        state_path = os.path.join(session_folder, 'extract_state.json')
        with open(state_path, 'w', encoding='utf-8') as f:
            json.dump(extract_state, f, ensure_ascii=False)

        yield _evt('writing', 88, message='Đang tính dedup...')

        # Dedup
        dedup_files, dedup_mapping, dedup_stats = build_dedup_data(extracted_data, CHUNK_SIZE)
        dedup_map_path = os.path.join(session_folder, 'dedup_mapping.json')
        with open(dedup_map_path, 'w', encoding='utf-8') as f:
            json.dump(dedup_mapping, f, ensure_ascii=False, indent=2)

        result = {
            'success': True,
            'total_files': num_files,
            'total_items': total_items,
            'files': files_data,
            'zip_display_name': zip_display,
            'dedup_files': dedup_files,
            'dedup_stats': dedup_stats,
        }

        yield _evt('done', 100, result=result, message='Hoàn tất!')

    except Exception as e:
        yield _evt('error', 0, error=f'Lỗi khi xử lý file: {str(e)}')


# ==================== SMART UPDATE (MATRIX MAPPING INHERITANCE) ====================

import tempfile as _tempfile
from openpyxl.cell import MergedCell as _MergedCell

# Mặc định: màu xanh lá trong VN_1.1 = nội dung mới/sửa
DEFAULT_NEW_COLORS = {'38761D', '00FF00', '008000', '006100', '00B050'}
DEFAULT_RED_COLORS  = {'FF0000', 'FF0000'}


def _get_rgb6(cell) -> str:
    """Đọc mã RGB 6 ký tự từ font.color, trả rỗng nếu không xác định được."""
    try:
        color = cell.font.color if (cell.font and cell.font.color) else None
        if color and color.type == 'rgb' and color.rgb:
            return str(color.rgb).upper().lstrip('F')[-6:].zfill(6)
    except Exception:
        pass
    return ''


# ==================== COLOR FILTER HELPERS ====================

def _get_font_rgb_xlsx(cell) -> str:
    """Đọc màu chữ HEX 6 ký tự từ openpyxl cell. Trả '' nếu không xác định."""
    try:
        color = cell.font.color if (cell.font and cell.font.color) else None
        if color and color.type == 'rgb' and color.rgb:
            return str(color.rgb).upper()[-6:]  # lấy 6 ký tự cuối của ARGB
    except Exception:
        pass
    return ''


def _get_font_rgb_pptx(run) -> str:
    """Đọc màu chữ HEX 6 ký tự từ python-pptx run. Trả '' nếu không xác định."""
    try:
        rgb = run.font.color.rgb
        if rgb is not None:
            return str(rgb).upper()
    except Exception:
        pass
    return ''


def _get_font_rgb_docx(run) -> str:
    """Đọc màu chữ HEX 6 ký tự từ python-docx run. Trả '' nếu không xác định."""
    try:
        rgb = run.font.color.rgb
        if rgb is not None:
            return str(rgb).upper()
    except Exception:
        pass
    return ''


def _normalize_color_filter(color_list: list) -> set:
    """Chuẩn hóa list màu HEX → set HEX 6 UPPERCASE không '#'."""
    result = set()
    for c in color_list:
        c = c.strip().lstrip('#').upper()
        if len(c) == 6:
            result.add(c)
    return result


_PROOF_URL_RE = re.compile(r'^(?:https?://|ftp://|www\.)\S+$', re.IGNORECASE)
_PROOF_EMAIL_RE = re.compile(r'^[^\s@]+@[^\s@]+\.[^\s@]+$')
_PROOF_DATE_RE = re.compile(r'^\d{1,4}[\-/]\d{1,2}[\-/]\d{1,4}$')
_PROOF_TIME_RE = re.compile(r'^\d{1,2}:\d{2}(?::\d{2})?$')
_PROOF_NUMERIC_RE = re.compile(r'^[\d\s.,:+\-/%$€¥₫()]+$')
_PROOF_TECH_TOKEN_RE = re.compile(r'^[A-Za-z0-9][A-Za-z0-9_.:/\\#?=&%+~\-]*$')


def _is_proofread_excluded_text(value: str) -> bool:
    """True khi value không phải nội dung cần AI check ngữ pháp."""
    s = str(value).strip()
    if not s:
        return True
    if _PROOF_URL_RE.match(s) or _PROOF_EMAIL_RE.match(s):
        return True
    if _PROOF_DATE_RE.match(s) or _PROOF_TIME_RE.match(s):
        return True
    if _PROOF_NUMERIC_RE.match(s):
        return True
    if not any(ch.isalpha() for ch in s):
        return True
    if ' ' not in s and _PROOF_TECH_TOKEN_RE.match(s):
        has_digit = any(ch.isdigit() for ch in s)
        has_tech_sep = any(ch in '._-:/\\#?=&%+' for ch in s)
        if has_digit or has_tech_sep:
            return True
    return False


def _filter_proofread_extract_data(extracted_data: dict) -> dict:
    """Lọc bỏ URL/chuỗi kỹ thuật/số-ngày-ký hiệu khỏi dữ liệu trích xuất cho Tab 5."""
    filtered = {}
    for key, value in extracted_data.items():
        if isinstance(value, str) and _is_proofread_excluded_text(value):
            continue
        filtered[key] = value
    return filtered


def _pptx_shape_matches_color_filter(shape, color_filter: set) -> bool:
    """
    Trả về True nếu bất kỳ run nào trong text_frame của shape khớp color_filter.
    Nếu không có run nào → match khi '000000' trong filter (màu mặc định).
    """
    if not hasattr(shape, 'text_frame'):
        return True  # không kiểm tra được → include mặc định
    all_runs = [run for para in shape.text_frame.paragraphs for run in para.runs]
    if not all_runs:
        return '000000' in color_filter
    return any((_get_font_rgb_pptx(run) or '000000') in color_filter for run in all_runs)


def _pptx_cell_matches_color_filter(cell, color_filter: set) -> bool:
    """Trả về True nếu bất kỳ run nào trong table cell của pptx khớp color_filter."""
    tf = getattr(cell, 'text_frame', None)
    if tf is None:
        return True
    all_runs = [run for para in tf.paragraphs for run in para.runs]
    if not all_runs:
        return '000000' in color_filter
    return any((_get_font_rgb_pptx(run) or '000000') in color_filter for run in all_runs)


def _docx_para_matches_color_filter(para, color_filter: set) -> bool:
    """Trả về True nếu bất kỳ run nào trong paragraph docx khớp color_filter."""
    if not para.runs:
        return '000000' in color_filter
    return any((_get_font_rgb_docx(run) or '000000') in color_filter for run in para.runs)


def _docx_cell_matches_color_filter(cell, color_filter: set) -> bool:
    """Trả về True nếu bất kỳ run nào trong table cell docx khớp color_filter."""
    has_any_run = False
    for para in cell.paragraphs:
        for run in para.runs:
            has_any_run = True
            if (_get_font_rgb_docx(run) or '000000') in color_filter:
                return True
    if not has_any_run:
        return '000000' in color_filter
    return False


def _collect_pptx_shape_colors(shape, colors: set):
    """Thu thập màu chữ từ tất cả runs trong shape pptx (đệ quy cho grouped shapes)."""
    if hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                c = _get_font_rgb_pptx(run)
                if c:
                    colors.add(c)
    if hasattr(shape, 'has_table') and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                tf = getattr(cell, 'text_frame', None)
                if tf:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            c = _get_font_rgb_pptx(run)
                            if c:
                                colors.add(c)
    if hasattr(shape, 'shapes'):
        for child_shape in shape.shapes:
            _collect_pptx_shape_colors(child_shape, colors)


def is_new_or_modified_cell(cell_vn11, new_colors=None, red_colors=None) -> bool:
    """
    Trả về True nếu ô được đánh dấu màu xanh lá (nội dung mới/sửa).
    Trả về False nếu màu đỏ (marker) hoặc đen/auto (nội dung cũ).
    None → không xác định qua màu, cần dng fallback.
    """
    if new_colors is None:
        new_colors = DEFAULT_NEW_COLORS
    if red_colors is None:
        red_colors = DEFAULT_RED_COLORS
    rgb = _get_rgb6(cell_vn11)
    if not rgb:
        return None  # không có màu rõ ràng → fallback
    if rgb in new_colors:
        return True
    if rgb in red_colors:
        return False  # marker, bỏ qua
    return None  # màu khác → fallback sang so sánh nội dung


def is_content_changed(cell_vn10, cell_vn11) -> bool:
    """Fallback: so sánh nội dung text giữa VN_1.0 và VN_1.1 tại cùng tọa độ."""
    v0 = str(cell_vn10.value).strip() if cell_vn10.value is not None else ''
    v1 = str(cell_vn11.value).strip() if cell_vn11.value is not None else ''
    return bool(v1) and v1 != v0


def _should_skip_cell(cell) -> bool:
    """Bỏ qua ô công thức, merged non-top-left, hoặc rỗng."""
    if isinstance(cell, _MergedCell):
        return True
    if cell.value is None:
        return True
    if isinstance(cell.value, str):
        if cell.value.strip() == '':
            return True
        if cell.value.startswith('='):
            return True
    return False


def _safe_set_value(ws, coord, value):
    """
    Ghi giá trị vào ô worksheet, bỏ qua nếu ô đó là MergedCell (non-top-left).
    Tránh lỗi ‘MergedCell’ object attribute ‘value’ is read-only.
    """
    cell = ws[coord]
    if isinstance(cell, _MergedCell):
        return  # chỉ top-left của merge mới ghi được
    cell.value = value


def _copy_cell_format(src_cell, dst_cell):
    """
    Copy định dạng (format) từ ô nguồn sang ô đích.
    Bao gồm: number_format, font, alignment, fill, border.
    Dùng để bảo toàn định dạng khi kế thừa giá trị dịch từ JP_1.0 sang JP_1.1.
    """
    if src_cell is None or dst_cell is None:
        return
    try:
        # Copy number format (định dạng ngày/số)
        if src_cell.number_format:
            dst_cell.number_format = src_cell.number_format
        
        # Copy font
        if src_cell.font:
            dst_cell.font = copy(src_cell.font)
        
        # Copy alignment
        if src_cell.alignment:
            dst_cell.alignment = copy(src_cell.alignment)
        
        # Copy fill (màu nền)
        if src_cell.fill:
            dst_cell.fill = copy(src_cell.fill)
        
        # Copy border
        if src_cell.border:
            dst_cell.border = copy(src_cell.border)
    except Exception:
        pass  # Bỏ qua nếu copy format thất bại


def _build_coord_content_map(ws_vn10, ws_jp10) -> dict:
    """
    Map: { (coordinate, vn_text) → jp_text }

    Key là TUPLE (tọa độ, nội dung VN) — chỉ kế thừa JP khi CẢ HAI đều khớp:
    - Tọa độ giống nhau (coord)
    - Nội dung VN tại tọa độ đó không thay đổi (vn_text)

    Đây là tầng chính xác nhất: nếu VN_1.1[A5] == VN_1.0[A5] thì JP_1.0[A5]
    chắc chắn là bản dịch đúng cho ô đó, bất kể text có trùng với ô khác không.
    """
    result = {}
    for row in ws_vn10.iter_rows():
        for cell_vn in row:
            if isinstance(cell_vn, _MergedCell) or not cell_vn.value:
                continue
            vn_text = str(cell_vn.value).strip()
            if not vn_text or vn_text.startswith('='):
                continue
            jp_cell = ws_jp10[cell_vn.coordinate]
            if isinstance(jp_cell, _MergedCell) or not jp_cell.value:
                continue
            jp_text = str(jp_cell.value).strip()
            if jp_text:
                result[(cell_vn.coordinate, vn_text)] = jp_text
    return result


def _build_vn_jp_content_map(ws_vn10, ws_jp10) -> dict:
    """
    Map: { vn_text → (jp_text, coord) } — fallback khi tọa độ đã thay đổi.

    Với mỗi VN text, đếm số lần xuất hiện của từng bản dịch JP tương ứng,
    rồi chọn JP xuất hiện NHIỀU NHẤT (dominant).

    Lý do dùng dominant thay vì loại bỏ conflict: bản dịch JP xuất hiện
    nhiều nhất là bản khách đã chấp nhận nhiều lần — đáng tin cậy hơn.

    Returns: {vn_text: (jp_text, coord_đầu_tiên_có_jp_text_này)}
    """
    from collections import defaultdict
    vn_to_jp_counter = defaultdict(lambda: defaultdict(int))
    vn_to_jp_coord = {}  # {vn_text: {jp_text: coord_đầu_tiên}}

    for row in ws_vn10.iter_rows():
        for cell_vn in row:
            if isinstance(cell_vn, _MergedCell) or not cell_vn.value:
                continue
            vn_text = str(cell_vn.value).strip()
            if not vn_text or vn_text.startswith('='):
                continue
            jp_cell = ws_jp10[cell_vn.coordinate]
            if isinstance(jp_cell, _MergedCell) or not jp_cell.value:
                continue
            jp_text = str(jp_cell.value).strip()
            if jp_text:
                vn_to_jp_counter[vn_text][jp_text] += 1
                # Lưu coordinate của jp_cell (lần đầu tiên gặp cặp này)
                if vn_text not in vn_to_jp_coord:
                    vn_to_jp_coord[vn_text] = {}
                if jp_text not in vn_to_jp_coord[vn_text]:
                    vn_to_jp_coord[vn_text][jp_text] = jp_cell.coordinate

    result = {}
    for vn_text, jp_counter in vn_to_jp_counter.items():
        dominant_jp = max(jp_counter, key=jp_counter.get)
        coord = vn_to_jp_coord[vn_text][dominant_jp]
        result[vn_text] = (dominant_jp, coord)  # (value, coordinate)
    
    return result


def _clone_vn11_as_base(path_vn11: str, tmp_path: str) -> None:
    """
    Clone VN_1.1 thành file base cho JP_1.1.
    Giữ nguyên TOÀN BỘ cấu trúc (merge cells, row heights, col widths, styles).
    Sau đó xóa sạch text trong các cell (giữ format) để điền JP vào.
    """
    import shutil
    shutil.copy2(path_vn11, tmp_path)
    wb = load_workbook(tmp_path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell, _MergedCell):
                    continue
                if cell.value is not None:
                    cell.value = None
    wb.save(tmp_path)
    wb.close()


def _copy_sheet_structure(ws_source, ws_dest):
    """Copy column widths, row heights, merge cells từ source sang dest."""
    for merge in ws_source.merged_cells.ranges:
        try:
            ws_dest.merge_cells(str(merge))
        except Exception:
            pass
    for col_letter, col_dim in ws_source.column_dimensions.items():
        ws_dest.column_dimensions[col_letter].width = col_dim.width
    for row_num, row_dim in ws_source.row_dimensions.items():
        ws_dest.row_dimensions[row_num].height = row_dim.height


def smart_update_excel(path_vn10, path_vn11, path_jp10, new_colors=None, red_colors=None):
    """
    Tạo JP_1.1 bằng cách:
      - Clone VN_1.1 làm base (cấu trúc đúng nhất: merge, rows, cols)
      - Ô màu đỏ (marker)       → giữ text VN gốc
      - Ô màu xanh (new/sửa)    → giữ text VN_1.1, thêm vào to_translate
      - Ô đen/auto               → tầng 1 (coord+text) → tầng 2 (dominant JP) → giữ text VN_1.1
      - Sheet mới trong VN_1.1   → toàn bộ giữ text VN_1.1 + to_translate
      - Sheet bỏ trong VN_1.1    → bỏ qua (JP_1.1 theo cấu trúc VN_1.1)
      - Sheet có trong VN nhưng JP_1.0 thiếu → dùng text VN_1.1 thức đẩy to_translate
    Returns: (wb_jp11, to_translate_dict, stats_dict)
    """
    import shutil

    wb_vn10 = load_workbook(path_vn10, data_only=True)
    wb_vn11 = load_workbook(path_vn11, data_only=True)
    wb_jp10 = load_workbook(path_jp10, data_only=True)

    # Clone VN_1.1 → base của JP_1.1 (cấu trúc đúng nhất)
    tmp_path = _tempfile.mktemp(suffix='_jp11.xlsx')
    _clone_vn11_as_base(path_vn11, tmp_path)
    wb_jp11 = load_workbook(tmp_path)

    to_translate = {}
    inherited = 0
    sheet_stats = {}   # per-sheet breakdown

    for sheet_name in wb_vn11.sheetnames:
        ws_vn11 = wb_vn11[sheet_name]
        ws_vn10 = wb_vn10[sheet_name] if sheet_name in wb_vn10.sheetnames else None
        ws_jp10 = wb_jp10[sheet_name] if sheet_name in wb_jp10.sheetnames else None
        ws_jp11 = wb_jp11[sheet_name] if sheet_name in wb_jp11.sheetnames \
            else wb_jp11.create_sheet(sheet_name)

        sh_inherited   = 0
        sh_to_translate = 0

        # Case A: Sheet mới trong VN_1.1 (không có trong VN_1.0 và/hoặc JP_1.0)
        #         → toàn bộ ô giữ text VN_1.1 và được đánh dấu cần dịch
        if ws_vn10 is None:
            for row in ws_vn11.iter_rows():
                for cell in row:
                    if _should_skip_cell(cell):
                        continue
                    key = f"{sheet_name}!{cell.coordinate}"
                    to_translate[key] = str(cell.value).strip()
                    _safe_set_value(ws_jp11, cell.coordinate, cell.value)  # giữ text VN_1.1
                    sh_to_translate += 1
            sheet_stats[sheet_name] = {'inherited': 0, 'to_translate': sh_to_translate, 'status': 'new_sheet'}
            continue

        # Case B: Sheet tồn tại trong cả VN_1.0 và VN_1.1 → xử lý kế thừa
        # ws_jp10 có thể là None (sheet chưa tồn tại trong JP_1.0) → map rỗng → tất cả đen→translate
        coord_content_map = _build_coord_content_map(ws_vn10, ws_jp10) \
            if (ws_vn10 and ws_jp10) else {}
        vn_jp_content_map = _build_vn_jp_content_map(ws_vn10, ws_jp10) \
            if (ws_vn10 and ws_jp10) else {}

        for row in ws_vn11.iter_rows():
            for cell_vn11 in row:
                if _should_skip_cell(cell_vn11):
                    continue

                coord     = cell_vn11.coordinate
                text_vn11 = str(cell_vn11.value).strip()
                key       = f"{sheet_name}!{coord}"

                # ── Non-string cells: int, float, datetime, bool ──
                # Không cần dịch, không cần tra maps — giữ nguyên native type.
                # Maps chỉ xử lý str; ghi str vào cell sẽ làm mất kiểu và format Excel.
                if not isinstance(cell_vn11.value, str):
                    _safe_set_value(ws_jp11, coord, cell_vn11.value)
                    inherited += 1
                    sh_inherited += 1
                    continue

                color_result = is_new_or_modified_cell(cell_vn11, new_colors, red_colors)

                if color_result is False:
                    # Màu đỏ = marker (①②...) → giữ nguyên text VN trong JP_1.1
                    _safe_set_value(ws_jp11, coord, cell_vn11.value)
                    continue

                if color_result is True:
                    # Màu xanh = nội dung mới/sửa → cần dịch mới
                    to_translate[key] = text_vn11
                    _safe_set_value(ws_jp11, coord, cell_vn11.value)  # giữ text VN_1.1
                    sh_to_translate += 1
                    continue

                # Không màu rõ ràng → fallback theo 2 tầng

                # Tầng 1: coord + vn_text đều khớp VN_1.0 → JP_1.0[coord] chắc chắn đúng
                jp_val = coord_content_map.get((coord, text_vn11))
                if jp_val:
                    _safe_set_value(ws_jp11, coord, jp_val)
                    # Copy định dạng từ JP_1.0 sang JP_1.1
                    src_cell = ws_jp10[coord]
                    dst_cell = ws_jp11[coord]
                    _copy_cell_format(src_cell, dst_cell)
                    inherited += 1
                    sh_inherited += 1
                    continue

                # Tầng 2: vn_text có trong VN_1.0 → lấy JP dominant (bản khách đã chấp nhận)
                if text_vn11 in vn_jp_content_map:
                    jp_val, jp_coord = vn_jp_content_map[text_vn11]
                    _safe_set_value(ws_jp11, coord, jp_val)
                    # Copy định dạng từ JP_1.0[jp_coord] sang JP_1.1[coord]
                    src_cell = ws_jp10[jp_coord]
                    dst_cell = ws_jp11[coord]
                    _copy_cell_format(src_cell, dst_cell)
                    inherited += 1
                    sh_inherited += 1
                    continue

                # Không tìm được → text mới hoàn toàn → cần dịch mới
                to_translate[key] = text_vn11
                _safe_set_value(ws_jp11, coord, cell_vn11.value)  # giữ text VN_1.1
                sh_to_translate += 1

        sheet_stats[sheet_name] = {
            'inherited': sh_inherited,
            'to_translate': sh_to_translate,
            'status': 'updated' if ws_jp10 else 'no_jp10',
        }

    wb_vn10.close()
    wb_vn11.close()
    wb_jp10.close()

    try:
        os.remove(tmp_path)
    except Exception:
        pass

    stats = {
        'inherited': inherited,
        'to_translate': len(to_translate),
        'total': inherited + len(to_translate),
        'sheets': sheet_stats,
    }
    return wb_jp11, to_translate, stats


# giữ alias để bảo tương thích ngược với code cũ (nếu có chỗ nào gọi)
compare_and_inherit_excel = smart_update_excel


def get_password():
    """Đọc password từ file password.txt"""
    try:
        with open(PASSWORD_FILE, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except FileNotFoundError:
        # Nếu file không tồn tại, tạo file với password mặc định
        default_password = 'admin123'
        with open(PASSWORD_FILE, 'w', encoding='utf-8') as f:
            f.write(default_password)
        return default_password

def get_machine_id():
    """Lấy ID máy (dựa trên UUID node)"""
    return hex(uuid.getnode())

def create_session_id():
    """Tạo session ID dựa trên machine ID + timestamp"""
    machine_id = get_machine_id()
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    return f"{machine_id}_{timestamp}"

def get_session_folder():
    """Lấy đường dẫn folder của session hiện tại"""
    if 'session_id' not in session:
        session['session_id'] = create_session_id()
    
    session_folder = os.path.join(app.config['UPLOAD_FOLDER'], session['session_id'])
    os.makedirs(session_folder, exist_ok=True)
    return session_folder

def cleanup_old_sessions():
    """Xóa tất cả folder của các phiên từ hôm qua trở về trước"""
    try:
        upload_folder = app.config['UPLOAD_FOLDER']
        if not os.path.exists(upload_folder):
            return
        
        # Lấy ngày hiện tại (không có giờ phút giây)
        today = datetime.now().date()
        
        # Duyệt qua tất cả các folder trong uploads
        for folder_name in os.listdir(upload_folder):
            folder_path = os.path.join(upload_folder, folder_name)
            
            if os.path.isdir(folder_path):
                try:
                    # Parse timestamp từ tên folder (format: machine_YYYYMMDD_HHMMSS)
                    parts = folder_name.split('_')
                    if len(parts) >= 2:
                        date_str = parts[-2]  # YYYYMMDD
                        folder_date = datetime.strptime(date_str, '%Y%m%d').date()
                        
                        # Nếu folder từ hôm qua trở về trước, xóa đi
                        if folder_date < today:
                            shutil.rmtree(folder_path)
                            print(f"Đã xóa folder cũ: {folder_name}")
                except (ValueError, IndexError):
                    # Nếu không parse được, bỏ qua
                    continue
    except Exception as e:
        print(f"Lỗi khi cleanup old sessions: {e}")

def login_required(f):
    """Decorator để yêu cầu đăng nhập"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            # Nếu là request AJAX/JSON, trả về JSON thay vì redirect
            if request.path.startswith('/api') or request.is_json or request.path in ['/extract', '/inject', '/clear-uploads']:
                return jsonify({'error': 'Chưa đăng nhập hoặc phiên đã hết hạn'}), 401
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def allowed_file(filename):
    """
    Kiểm tra xem file có phải định dạng được phép không
    """
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def set_download_headers(response, display_name, default_ascii_name):
    """Set Content-Disposition hỗ trợ tên file Unicode (RFC 5987)."""
    encoded_filename = quote(display_name, safe='')
    ascii_filename = secure_filename(display_name) or default_ascii_name
    response.headers['Content-Disposition'] = (
        f"attachment; filename=\"{ascii_filename}\"; filename*=UTF-8''{encoded_filename}"
    )
    return response

# Error Handlers
@app.errorhandler(400)
def bad_request(error):
    """Xử lý lỗi 400 Bad Request"""
    if request.path.startswith('/api') or request.is_json or request.path in ['/extract', '/inject', '/clear-uploads']:
        return jsonify({'error': str(error) or 'Yêu cầu không hợp lệ'}), 400
    return str(error), 400

@app.errorhandler(401)
def unauthorized(error):
    """Xử lý lỗi 401 Unauthorized"""
    if request.path.startswith('/api') or request.is_json or request.path in ['/extract', '/inject', '/clear-uploads']:
        return jsonify({'error': 'Chưa đăng nhập hoặc phiên đã hết hạn'}), 401
    return redirect(url_for('login'))

@app.errorhandler(404)
def not_found(error):
    """Xử lý lỗi 404 Not Found"""
    if request.path.startswith('/api') or request.is_json:
        return jsonify({'error': 'Không tìm thấy tài nguyên'}), 404
    return str(error), 404

@app.errorhandler(413)
def request_entity_too_large(error):
    """Xử lý lỗi 413 Request Entity Too Large"""
    return jsonify({'error': 'File quá lớn. Giới hạn 50MB'}), 413

@app.errorhandler(500)
def internal_server_error(error):
    """Xử lý lỗi 500 Internal Server Error"""
    if request.path.startswith('/api') or request.is_json or request.path in ['/extract', '/inject', '/clear-uploads']:
        return jsonify({'error': f'Lỗi máy chủ: {str(error)}'}), 500
    return str(error), 500

def extract_text_from_shape(shape, shape_path, extracted_data, color_filter=None):
    """
    Hàm đệ quy để trích xuất text từ shape, bao gồm cả grouped shapes
    shape_path: đường dẫn đến shape, ví dụ "Shape1" hoặc "Shape1_2_3"
    color_filter: set HEX strings hoặc None (không lọc)
    """
    # Trích xuất text từ text frame của shape hiện tại
    if hasattr(shape, "text") and shape.text:
        text_content = shape.text.strip()
        if text_content:  # Chỉ lấy nội dung không rỗng
            if color_filter is None or _pptx_shape_matches_color_filter(shape, color_filter):
                extracted_data[shape_path] = text_content
    
    # Trích xuất text từ table nếu có
    if hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows, start=1):
            for col_idx, cell in enumerate(row.cells, start=1):
                if cell.text.strip():
                    if color_filter is None or _pptx_cell_matches_color_filter(cell, color_filter):
                        key = f"{shape_path}!Table_R{row_idx}C{col_idx}"
                        extracted_data[key] = cell.text.strip()
    
    # Kiểm tra xem shape có phải là GroupShape không (chứa các shape con)
    if hasattr(shape, "shapes"):
        # Đây là grouped shape, duyệt qua các shape con
        for child_idx, child_shape in enumerate(shape.shapes, start=1):
            child_path = f"{shape_path}_{child_idx}"
            extract_text_from_shape(child_shape, child_path, extracted_data, color_filter)

def extract_text_from_pptx(filepath, color_filter=None):
    """
    Trích xuất text từ file PPTX, bao gồm cả text trong grouped shapes
    Trả về dictionary với format: {"SlideX!ShapeY": "Content"}
    Với nested shapes: {"SlideX!ShapeY_Z": "Content"} (Z là shape con)
    color_filter: set HEX strings hoặc None (không lọc)
    """
    extracted_data = {}
    prs = Presentation(filepath)
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            shape_path = f"Slide{slide_idx}!Shape{shape_idx}"
            extract_text_from_shape(shape, shape_path, extracted_data, color_filter)
    
    return extracted_data

def inject_text_to_shape(shape, shape_indices, translated_value, is_table_cell=False, table_pos=None):
    """
    Hàm đệ quy để nạp text vào shape, bao gồm cả grouped shapes
    shape_indices: list các index để navigate đến shape đúng, ví dụ [2, 3] cho Shape2_3
    is_table_cell: có phải là table cell không
    table_pos: tuple (row_idx, col_idx) nếu là table cell
    """
    # Nếu là shape cuối cùng trong path
    if len(shape_indices) == 0:
        if is_table_cell and table_pos:
            # Nạp vào table cell - giữ nguyên định dạng
            row_idx, col_idx = table_pos
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                    cell = table.rows[row_idx].cells[col_idx]
                    # Thay thế text trong từng paragraph/run để giữ định dạng
                    if cell.text_frame:
                        replace_text_keep_format(cell.text_frame, translated_value)
        else:
            # Nạp vào text frame của shape - giữ nguyên định dạng
            if hasattr(shape, "text_frame") and shape.text_frame:
                replace_text_keep_format(shape.text_frame, translated_value)
        return True
    
    # Navigate đến shape con
    if hasattr(shape, "shapes"):
        next_idx = shape_indices[0]
        if next_idx <= len(shape.shapes):
            child_shape = shape.shapes[next_idx - 1]  # Chuyển từ 1-indexed sang 0-indexed
            return inject_text_to_shape(child_shape, shape_indices[1:], translated_value, is_table_cell, table_pos)
    
    return False

def replace_text_keep_format(text_frame, new_text):
    """
    Thay thế text trong text_frame nhưng giữ nguyên định dạng (font, màu, gạch chân, bold, italic...)
    Chiến lược:
    1. Nếu toàn bộ text frame chỉ có 1 paragraph và 1 run -> thay text của run đó
    2. Nếu có nhiều runs/paragraphs -> xóa text của tất cả runs, gán text mới vào run đầu tiên với định dạng gốc
    """
    if not text_frame.paragraphs:
        return
    
    # Thu thập tất cả runs từ tất cả paragraphs
    all_runs = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            all_runs.append(run)
    
    if not all_runs:
        # Không có run nào, tạo mới
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = new_text
        return
    
    # Lưu định dạng của run đầu tiên
    first_run = all_runs[0]
    
    # Xóa text của tất cả runs
    for run in all_runs:
        run.text = ""
    
    # Gán text mới vào run đầu tiên (giữ nguyên định dạng)
    first_run.text = new_text

def inject_text_to_pptx(filepath, json_data):
    """
    Nạp text đã dịch vào file PPTX, bao gồm cả grouped shapes
    """
    prs = Presentation(filepath)
    
    for key, translated_value in json_data.items():
        try:
            # Parse key format: 
            # "SlideX!ShapeY" hoặc "SlideX!ShapeY_Z" (nested) 
            # hoặc "SlideX!ShapeY!Table_RxCy" hoặc "SlideX!ShapeY_Z!Table_RxCy"
            if '!' not in key:
                continue
            
            parts = key.split('!')
            if len(parts) < 2:
                continue
            
            # Lấy slide index
            slide_part = parts[0]
            if not slide_part.startswith('Slide'):
                continue
            slide_idx = int(slide_part.replace('Slide', '')) - 1
            
            if slide_idx >= len(prs.slides):
                continue
            
            slide = prs.slides[slide_idx]
            
            # Parse shape path: "Shape2" hoặc "Shape2_3_1" (nested)
            shape_part = parts[1]
            if not shape_part.startswith('Shape'):
                continue
            
            # Tách các indices: "Shape2_3_1" -> [2, 3, 1]
            shape_str = shape_part.replace('Shape', '')
            shape_indices = [int(idx) for idx in shape_str.split('_')]
            
            # Lấy shape đầu tiên (top-level shape)
            first_shape_idx = shape_indices[0] - 1  # Chuyển sang 0-indexed
            if first_shape_idx >= len(slide.shapes):
                continue
            
            shape = slide.shapes[first_shape_idx]
            
            # Kiểm tra xem có phải table cell không
            is_table_cell = False
            table_pos = None
            
            if len(parts) == 3 and parts[2].startswith('Table_R'):
                # Parse table cell position
                is_table_cell = True
                table_part = parts[2].replace('Table_R', '').split('C')
                row_idx = int(table_part[0]) - 1
                col_idx = int(table_part[1]) - 1
                table_pos = (row_idx, col_idx)
            
            # Navigate và nạp text (bỏ qua index đầu tiên vì đã lấy shape rồi)
            inject_text_to_shape(shape, shape_indices[1:], translated_value, is_table_cell, table_pos)
            
        except (ValueError, IndexError, AttributeError) as e:
            # Bỏ qua các key không hợp lệ
            continue
    
    return prs

# ==================== XLSX SHAPE / OBJECT SUPPORT ====================
# Namespaces dùng trong drawing XML của xlsx
_NS_XDR = 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing'
_NS_A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_R   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
_NS_WB  = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'


def _xlsx_get_sheet_drawing_map(z):
    """
    Từ ZipFile đang mở, trả về dict: sheet_name → list các đường dẫn drawing XML
    Ví dụ: {'Sheet1': ['xl/drawings/drawing1.xml']}
    """
    names_set = set(z.namelist())

    # Đọc workbook.xml để lấy tên sheet và rId
    wb_xml = z.read('xl/workbook.xml')
    wb_root = _etree.fromstring(wb_xml)

    # Đọc rels của workbook để map rId → target file
    wb_rels_xml = z.read('xl/_rels/workbook.xml.rels')
    wb_rels_root = _etree.fromstring(wb_rels_xml)
    rid_to_target = {rel.get('Id'): rel.get('Target') for rel in wb_rels_root}

    # Thu thập (sheet_name, sheet_path)
    sheet_info = []
    for sheet_el in wb_root.iter(f'{{{_NS_WB}}}sheet'):
        name = sheet_el.get('name')
        rid  = sheet_el.get(f'{{{_NS_R}}}id')
        target = rid_to_target.get(rid, '')
        # Chuẩn hóa path: "worksheets/sheet1.xml" → "xl/worksheets/sheet1.xml"
        if target.startswith('../'):
            sheet_path = 'xl/' + target[3:]
        elif not target.startswith('xl/'):
            sheet_path = 'xl/' + target
        else:
            sheet_path = target
        sheet_info.append((name, sheet_path))

    result = {}
    for sheet_name, sheet_path in sheet_info:
        if sheet_path not in names_set:
            continue
        parts = sheet_path.rsplit('/', 1)
        sheet_dir  = parts[0] if len(parts) > 1 else ''
        sheet_file = parts[-1]
        rels_path  = f"{sheet_dir}/_rels/{sheet_file}.rels"
        if rels_path not in names_set:
            continue

        rels_xml  = z.read(rels_path)
        rels_root = _etree.fromstring(rels_xml)

        drawing_paths = []
        for rel in rels_root:
            if 'drawing' in rel.get('Type', '').lower():
                tgt = rel.get('Target', '')
                if tgt.startswith('../'):
                    draw_path = 'xl/' + tgt[3:]
                elif tgt.startswith('/'):
                    draw_path = tgt.lstrip('/')
                else:
                    draw_path = f"{sheet_dir}/{tgt.lstrip('./')}"
                if draw_path in names_set:
                    drawing_paths.append(draw_path)
        if drawing_paths:
            result[sheet_name] = drawing_paths

    return result


def _collect_sp_elements(drawing_root):
    """
    Thu thập tất cả phần tử <xdr:sp> (shape/text-box) theo thứ tự cây (tree-order),
    bao gồm cả sp bên trong group-shapes (xdr:grpSp).
    """
    sp_list = []
    _walk_tags = {
        f'{{{_NS_XDR}}}grpSp',
        f'{{{_NS_XDR}}}twoCellAnchor',
        f'{{{_NS_XDR}}}oneCellAnchor',
        f'{{{_NS_XDR}}}absoluteAnchor',
    }

    def walk(el):
        for child in el:
            if child.tag == f'{{{_NS_XDR}}}sp':
                sp_list.append(child)
            elif child.tag in _walk_tags:
                walk(child)

    walk(drawing_root)
    return sp_list


def _get_sp_text(sp):
    """Lấy toàn bộ text trong txBody của một shape."""
    txBody = sp.find(f'{{{_NS_XDR}}}txBody')
    if txBody is None:
        return ''
    parts = []
    for para in txBody.findall(f'{{{_NS_A}}}p'):
        # Lấy tất cả a:t trong paragraph (bao gồm cả text trong a:fld)
        para_text = ''.join((t.text or '') for t in para.findall(f'.//{{{_NS_A}}}t'))
        if para_text:
            parts.append(para_text)
    return '\n'.join(parts)


def _set_sp_text(sp, new_text):
    """
    Thay thế text trong txBody của shape theo cách tối thiểu để tương thích Excel:
    - Giữ nguyên cấu trúc paragraph/run/fld hiện có
    - Chỉ thay text ở node a:t đầu tiên, các node a:t còn lại set rỗng
    - Nếu chưa có a:t thì tạo mới tối thiểu trong paragraph đầu
    """
    txBody = sp.find(f'{{{_NS_XDR}}}txBody')
    if txBody is None:
        return

    paras = txBody.findall(f'{{{_NS_A}}}p')
    if not paras:
        return

    all_t_nodes = txBody.findall(f'.//{{{_NS_A}}}t')
    new_text = '' if new_text is None else str(new_text)

    if all_t_nodes:
        first_t = all_t_nodes[0]
        first_t.text = new_text
        if new_text != new_text.strip() or '\n' in new_text:
            first_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
        else:
            first_t.attrib.pop('{http://www.w3.org/XML/1998/namespace}space', None)

        for t in all_t_nodes[1:]:
            t.text = ''
        return

    # Không có a:t nào, tạo tối thiểu trong paragraph đầu
    first_para = paras[0]
    first_run = _etree.SubElement(first_para, f'{{{_NS_A}}}r')
    first_t = _etree.SubElement(first_run, f'{{{_NS_A}}}t')
    first_t.text = new_text
    if new_text != new_text.strip() or '\n' in new_text:
        first_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')


def extract_xlsx_shapes(filepath):
    """
    Trích xuất text từ tất cả shape/object (text-box) trong file xlsx.
    Trả về dict: {"SheetName!XLShape{n}": "text"}
    - n là thứ tự shape (đếm TẤT CẢ sp, bao gồm cả sp không có text) → index ổn định.
    """
    shapes_data = {}
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            sheet_drawing_map = _xlsx_get_sheet_drawing_map(z)
            for sheet_name, drawing_paths in sheet_drawing_map.items():
                for drawing_path in drawing_paths:
                    drawing_xml  = z.read(drawing_path)
                    drawing_root = _etree.fromstring(drawing_xml)
                    for shape_idx, sp in enumerate(_collect_sp_elements(drawing_root), start=1):
                        text = _get_sp_text(sp)
                        if text.strip():
                            key = f"{sheet_name}!XLShape{shape_idx}"
                            shapes_data[key] = text.strip()
    except Exception as e:
        print(f"Warning: Không thể trích xuất shapes từ xlsx: {e}")
    return shapes_data


def _xlsx_sheet_path_map(files):
    """Trả về map: sheet_name -> sheet_xml_path từ workbook + workbook.rels."""
    workbook_root = _etree.fromstring(files['xl/workbook.xml'])
    rels_root = _etree.fromstring(files['xl/_rels/workbook.xml.rels'])
    rid_to_target = {rel.get('Id'): rel.get('Target') for rel in rels_root}

    sheet_map = {}
    for sheet_el in workbook_root.iter(f'{{{_NS_WB}}}sheet'):
        sheet_name = sheet_el.get('name')
        rid = sheet_el.get(f'{{{_NS_R}}}id')
        target = rid_to_target.get(rid, '')
        if target.startswith('/'):
            path = target.lstrip('/')
        elif target.startswith('xl/'):
            path = target
        else:
            path = f'xl/{target}'
        sheet_map[sheet_name] = path
    return sheet_map


def _xlsx_sheet_drawing_map_from_files(files, sheet_map):
    """Trả về map: sheet_name -> [drawing_xml_paths]."""
    result = {}
    names_set = set(files.keys())

    for sheet_name, sheet_path in sheet_map.items():
        parts = sheet_path.rsplit('/', 1)
        sheet_dir = parts[0] if len(parts) > 1 else ''
        sheet_file = parts[-1]
        rels_path = f"{sheet_dir}/_rels/{sheet_file}.rels"
        if rels_path not in names_set:
            continue

        rels_root = _etree.fromstring(files[rels_path])
        drawing_paths = []
        for rel in rels_root:
            if 'drawing' in (rel.get('Type') or '').lower():
                tgt = rel.get('Target', '')
                if tgt.startswith('/'):
                    draw_path = tgt.lstrip('/')
                else:
                    draw_path = os.path.normpath(os.path.join(sheet_dir, tgt)).replace('\\', '/')
                if draw_path in names_set:
                    drawing_paths.append(draw_path)

        if drawing_paths:
            result[sheet_name] = drawing_paths

    return result


def _xlsx_set_cell_inline_text(sheet_root, cell_ref, text_value):
    """Gán text vào một cell theo kiểu inlineStr, giữ nguyên style của cell nếu có."""
    m = re.fullmatch(r'([A-Za-z]+)(\d+)', cell_ref or '')
    if not m:
        return

    row_num = int(m.group(2))
    sheet_data = sheet_root.find(f'{{{_NS_WB}}}sheetData')
    if sheet_data is None:
        return

    row_elem = None
    rows = sheet_data.findall(f'{{{_NS_WB}}}row')
    for row in rows:
        if int(row.get('r', '0') or 0) == row_num:
            row_elem = row
            break

    if row_elem is None:
        row_elem = _etree.Element(f'{{{_NS_WB}}}row', r=str(row_num))
        inserted = False
        for idx, row in enumerate(rows):
            existing = int(row.get('r', '0') or 0)
            if existing > row_num:
                sheet_data.insert(idx, row_elem)
                inserted = True
                break
        if not inserted:
            sheet_data.append(row_elem)

    cell_elem = None
    for cell in row_elem.findall(f'{{{_NS_WB}}}c'):
        if (cell.get('r') or '').upper() == cell_ref.upper():
            cell_elem = cell
            break

    if cell_elem is None:
        cell_elem = _etree.Element(f'{{{_NS_WB}}}c', r=cell_ref.upper())
        row_elem.append(cell_elem)

    text_str = '' if text_value is None else str(text_value)

    # Check if cell already has an <is> with rich-text runs
    existing_is = cell_elem.find(f'{{{_NS_WB}}}is')
    r_tag = f'{{{_NS_WB}}}r'
    t_tag = f'{{{_NS_WB}}}t'
    xml_space_attr = '{http://www.w3.org/XML/1998/namespace}space'

    if existing_is is not None:
        runs = existing_is.findall(r_tag)
        if len(runs) > 1:
            # Rich text: distribute new_text across runs by original char-count ratio
            orig_lengths = [len((r.findtext(t_tag) or '')) for r in runs]
            total_orig = sum(orig_lengths)
            new_total = len(text_str)
            pos = 0
            for i, (run_elem, orig_len) in enumerate(zip(runs, orig_lengths)):
                t_elem = run_elem.find(t_tag)
                if t_elem is None:
                    t_elem = _etree.SubElement(run_elem, t_tag)
                if i == len(runs) - 1:
                    part = text_str[pos:]
                else:
                    if total_orig > 0:
                        count = round(new_total * orig_len / total_orig)
                        count = min(count, new_total - pos - (len(runs) - i - 1))
                        count = max(count, 0)
                    else:
                        count = 0
                    part = text_str[pos:pos + count]
                    pos += count
                t_elem.text = part
                if part != part.strip() or '\n' in part:
                    t_elem.set(xml_space_attr, 'preserve')
                elif xml_space_attr in t_elem.attrib:
                    del t_elem.attrib[xml_space_attr]
            # Ensure cell type is inlineStr
            cell_elem.set('t', 'inlineStr')
            for child in list(cell_elem):
                if child.tag in {f'{{{_NS_WB}}}v', f'{{{_NS_WB}}}f'}:
                    cell_elem.remove(child)
            return
        elif len(runs) == 1:
            # Single run: just replace the text element, keep rPr
            t_elem = runs[0].find(t_tag)
            if t_elem is None:
                t_elem = _etree.SubElement(runs[0], t_tag)
            t_elem.text = text_str
            if text_str != text_str.strip() or '\n' in text_str:
                t_elem.set(xml_space_attr, 'preserve')
            elif xml_space_attr in t_elem.attrib:
                del t_elem.attrib[xml_space_attr]
            cell_elem.set('t', 'inlineStr')
            for child in list(cell_elem):
                if child.tag in {f'{{{_NS_WB}}}v', f'{{{_NS_WB}}}f'}:
                    cell_elem.remove(child)
            return
        # else: plain <is><t>...</t></is> — fall through to rewrite below

    # No existing <is>, or plain <is><t> only: rewrite from scratch
    cell_elem.set('t', 'inlineStr')
    for child in list(cell_elem):
        if child.tag in {
            f'{{{_NS_WB}}}v',
            f'{{{_NS_WB}}}f',
            f'{{{_NS_WB}}}is',
        }:
            cell_elem.remove(child)

    is_elem = _etree.SubElement(cell_elem, f'{{{_NS_WB}}}is')
    t_elem = _etree.SubElement(is_elem, t_tag)
    t_elem.text = text_str
    if text_str != text_str.strip() or '\n' in text_str:
        t_elem.set(xml_space_attr, 'preserve')


def inject_xlsx_shapes(source_filepath, output_filepath, json_data):
    """
    ZIP-level patch cho XLSX:
    - Không dùng openpyxl.save
    - Patch trực tiếp cell XML
    - Patch trực tiếp drawing XML (shape text)
    - Giữ nguyên toàn bộ parts khác của file gốc
    """
    with zipfile.ZipFile(source_filepath, 'r') as z_src:
        infos = {info.filename: info for info in z_src.infolist()}
        order = [info.filename for info in z_src.infolist()]
        files = {name: z_src.read(name) for name in order}

    sheet_map = _xlsx_sheet_path_map(files)

    cell_updates = {}
    shape_updates = {}
    for key, translated_value in json_data.items():
        if '!' not in key:
            continue
        sheet_name, second_part = key.split('!', 1)
        if sheet_name not in sheet_map:
            continue

        if second_part.startswith('XLShape'):
            try:
                shape_idx = int(second_part.replace('XLShape', ''))
                shape_updates[(sheet_name, shape_idx)] = '' if translated_value is None else str(translated_value)
            except ValueError:
                continue
        else:
            cell_updates.setdefault(sheet_name, {})[second_part] = translated_value

    for sheet_name, updates in cell_updates.items():
        sheet_path = sheet_map.get(sheet_name)
        if not sheet_path or sheet_path not in files:
            continue
        sheet_root = _etree.fromstring(files[sheet_path])
        for cell_ref, value in updates.items():
            _xlsx_set_cell_inline_text(sheet_root, cell_ref, value)
        # Sync hyperlink display attributes to match updated cell values
        hyperlinks_elem = sheet_root.find(f'{{{_NS_WB}}}hyperlinks')
        if hyperlinks_elem is not None:
            for hl in hyperlinks_elem.findall(f'{{{_NS_WB}}}hyperlink'):
                ref = hl.get('ref', '')
                if ref in updates and hl.get('display') is not None:
                    hl.set('display', str(updates[ref]))
        files[sheet_path] = _etree.tostring(sheet_root, xml_declaration=True, encoding='UTF-8', standalone=True)

    if shape_updates:
        drawing_map = _xlsx_sheet_drawing_map_from_files(files, sheet_map)
        for sheet_name, drawing_paths in drawing_map.items():
            wanted = {
                idx: val
                for (sn, idx), val in shape_updates.items()
                if sn == sheet_name
            }
            if not wanted:
                continue

            for drawing_path in drawing_paths:
                if drawing_path not in files:
                    continue
                drawing_root = _etree.fromstring(files[drawing_path])
                for shape_idx, sp in enumerate(_collect_sp_elements(drawing_root), start=1):
                    if shape_idx in wanted:
                        _set_sp_text(sp, wanted[shape_idx])
                files[drawing_path] = _etree.tostring(
                    drawing_root,
                    xml_declaration=True,
                    encoding='UTF-8',
                    standalone=True,
                )

    with zipfile.ZipFile(output_filepath, 'w') as z_out:
        written = set()
        for name in order:
            content = files.get(name)
            if content is None:
                continue
            original = infos[name]
            zi = zipfile.ZipInfo(name, original.date_time)
            zi.compress_type = original.compress_type
            zi.external_attr = original.external_attr
            zi.create_system = original.create_system
            z_out.writestr(zi, content)
            written.add(name)

        for name, content in files.items():
            if name in written:
                continue
            z_out.writestr(name, content)


def extract_text_from_docx(filepath, color_filter=None):
    """
    Trích xuất text từ file DOCX, bao gồm paragraphs, tables, headers, footers
    Trả về dictionary với format:
    - Paragraphs: {"ParagraphX": "Content"}
    - Tables: {"TableX!RyCz": "Content"}
    - Headers: {"Header_SectionX!ParagraphY": "Content"}
    - Footers: {"Footer_SectionX!ParagraphY": "Content"}
    color_filter: set HEX strings hoặc None (không lọc)
    """
    extracted_data = {}
    doc = Document(filepath)
    
    # 1. Trích xuất text từ các paragraph thông thường (không trong table)
    # Luôn đếm TẤT CẢ paragraph không rỗng để giữ index nhất quán với inject
    paragraph_idx = 0
    for para in doc.paragraphs:
        text_content = para.text.strip()
        if text_content:  # Đếm tất cả paragraph không rỗng
            paragraph_idx += 1
            if color_filter is None or _docx_para_matches_color_filter(para, color_filter):
                key = f"Paragraph{paragraph_idx}"
                extracted_data[key] = text_content
    
    # 2. Trích xuất text từ các bảng
    for table_idx, table in enumerate(doc.tables, start=1):
        for row_idx, row in enumerate(table.rows, start=1):
            for col_idx, cell in enumerate(row.cells, start=1):
                text_content = cell.text.strip()
                if text_content:
                    if color_filter is None or _docx_cell_matches_color_filter(cell, color_filter):
                        key = f"Table{table_idx}!R{row_idx}C{col_idx}"
                        extracted_data[key] = text_content
    
    # 3. Trích xuất text từ headers
    for section_idx, section in enumerate(doc.sections, start=1):
        header = section.header
        for para_idx, para in enumerate(header.paragraphs, start=1):
            text_content = para.text.strip()
            if text_content:
                if color_filter is None or _docx_para_matches_color_filter(para, color_filter):
                    key = f"Header_Section{section_idx}!Paragraph{para_idx}"
                    extracted_data[key] = text_content
        
        # Trích xuất từ table trong header (nếu có)
        for table_idx, table in enumerate(header.tables, start=1):
            for row_idx, row in enumerate(table.rows, start=1):
                for col_idx, cell in enumerate(row.cells, start=1):
                    text_content = cell.text.strip()
                    if text_content:
                        if color_filter is None or _docx_cell_matches_color_filter(cell, color_filter):
                            key = f"Header_Section{section_idx}!Table{table_idx}!R{row_idx}C{col_idx}"
                            extracted_data[key] = text_content
    
    # 4. Trích xuất text từ footers
    for section_idx, section in enumerate(doc.sections, start=1):
        footer = section.footer
        for para_idx, para in enumerate(footer.paragraphs, start=1):
            text_content = para.text.strip()
            if text_content:
                if color_filter is None or _docx_para_matches_color_filter(para, color_filter):
                    key = f"Footer_Section{section_idx}!Paragraph{para_idx}"
                    extracted_data[key] = text_content
        
        # Trích xuất từ table trong footer (nếu có)
        for table_idx, table in enumerate(footer.tables, start=1):
            for row_idx, row in enumerate(table.rows, start=1):
                for col_idx, cell in enumerate(row.cells, start=1):
                    text_content = cell.text.strip()
                    if text_content:
                        if color_filter is None or _docx_cell_matches_color_filter(cell, color_filter):
                            key = f"Footer_Section{section_idx}!Table{table_idx}!R{row_idx}C{col_idx}"
                            extracted_data[key] = text_content
    
    return extracted_data

def replace_text_keep_format_docx(paragraph, new_text):
    """
    Thay thế text trong paragraph nhưng GIỮ NGUYÊN format của từng run riêng lẻ.
    Chiến lược: phân phối new_text vào các runs theo tỉ lệ ký tự gốc.
    """
    if not paragraph.runs:
        paragraph.text = new_text
        return

    runs = paragraph.runs
    orig_lengths = [len(run.text) for run in runs]
    total_orig = sum(orig_lengths)

    if total_orig == 0:
        # All runs empty — assign everything to first run
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
        return

    new_total = len(new_text)
    pos = 0
    for i, (run, orig_len) in enumerate(zip(runs, orig_lengths)):
        if i == len(runs) - 1:
            run.text = new_text[pos:]
        else:
            count = round(new_total * orig_len / total_orig)
            count = min(count, new_total - pos - (len(runs) - i - 1))
            count = max(count, 0)
            run.text = new_text[pos:pos + count]
            pos += count

def inject_text_to_docx(filepath, json_data):
    """
    Nạp text đã dịch vào file DOCX
    Giữ nguyên định dạng (font, màu, size, bold, italic...)
    """
    doc = Document(filepath)
    
    for key, translated_value in json_data.items():
        try:
            # 1. Xử lý Paragraph thông thường: "ParagraphX"
            if key.startswith('Paragraph') and '!' not in key:
                para_num = int(key.replace('Paragraph', ''))
                # Đếm lại các paragraph không rỗng để map đúng index
                current_para_idx = 0
                for para in doc.paragraphs:
                    if para.text.strip():  # Chỉ đếm paragraph không rỗng
                        current_para_idx += 1
                        if current_para_idx == para_num:
                            replace_text_keep_format_docx(para, translated_value)
                            break
            
            # 2. Xử lý Table: "TableX!RyCz"
            elif key.startswith('Table') and '!' in key and not key.startswith('Header_') and not key.startswith('Footer_'):
                parts = key.split('!')
                if len(parts) != 2:
                    continue
                
                # Parse table index
                table_part = parts[0]
                table_idx = int(table_part.replace('Table', '')) - 1
                
                if table_idx >= len(doc.tables):
                    continue
                
                table = doc.tables[table_idx]
                
                # Parse cell position: "R2C3"
                cell_part = parts[1]
                if not cell_part.startswith('R'):
                    continue
                
                cell_parts = cell_part.replace('R', '').split('C')
                if len(cell_parts) != 2:
                    continue
                
                row_idx = int(cell_parts[0]) - 1
                col_idx = int(cell_parts[1]) - 1
                
                if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                    cell = table.rows[row_idx].cells[col_idx]
                    # Thay thế text trong paragraph đầu tiên của cell
                    if cell.paragraphs:
                        replace_text_keep_format_docx(cell.paragraphs[0], translated_value)
            
            # 3. Xử lý Header: "Header_SectionX!ParagraphY" hoặc "Header_SectionX!TableY!RzCw"
            elif key.startswith('Header_Section'):
                parts = key.split('!')
                if len(parts) < 2:
                    continue
                
                # Parse section index
                section_part = parts[0].replace('Header_Section', '')
                section_idx = int(section_part) - 1
                
                if section_idx >= len(doc.sections):
                    continue
                
                header = doc.sections[section_idx].header
                
                # Check if it's a paragraph or table
                if parts[1].startswith('Paragraph'):
                    para_num = int(parts[1].replace('Paragraph', ''))
                    para_idx = 0
                    for para in header.paragraphs:
                        if para.text.strip():
                            para_idx += 1
                            if para_idx == para_num:
                                replace_text_keep_format_docx(para, translated_value)
                                break
                
                elif parts[1].startswith('Table') and len(parts) == 3:
                    # Header table cell
                    table_idx = int(parts[1].replace('Table', '')) - 1
                    if table_idx >= len(header.tables):
                        continue
                    
                    table = header.tables[table_idx]
                    cell_part = parts[2]
                    if not cell_part.startswith('R'):
                        continue
                    
                    cell_parts = cell_part.replace('R', '').split('C')
                    if len(cell_parts) != 2:
                        continue
                    
                    row_idx = int(cell_parts[0]) - 1
                    col_idx = int(cell_parts[1]) - 1
                    
                    if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                        cell = table.rows[row_idx].cells[col_idx]
                        if cell.paragraphs:
                            replace_text_keep_format_docx(cell.paragraphs[0], translated_value)
            
            # 4. Xử lý Footer: "Footer_SectionX!ParagraphY" hoặc "Footer_SectionX!TableY!RzCw"
            elif key.startswith('Footer_Section'):
                parts = key.split('!')
                if len(parts) < 2:
                    continue
                
                # Parse section index
                section_part = parts[0].replace('Footer_Section', '')
                section_idx = int(section_part) - 1
                
                if section_idx >= len(doc.sections):
                    continue
                
                footer = doc.sections[section_idx].footer
                
                # Check if it's a paragraph or table
                if parts[1].startswith('Paragraph'):
                    para_num = int(parts[1].replace('Paragraph', ''))
                    para_idx = 0
                    for para in footer.paragraphs:
                        if para.text.strip():
                            para_idx += 1
                            if para_idx == para_num:
                                replace_text_keep_format_docx(para, translated_value)
                                break
                
                elif parts[1].startswith('Table') and len(parts) == 3:
                    # Footer table cell
                    table_idx = int(parts[1].replace('Table', '')) - 1
                    if table_idx >= len(footer.tables):
                        continue
                    
                    table = footer.tables[table_idx]
                    cell_part = parts[2]
                    if not cell_part.startswith('R'):
                        continue
                    
                    cell_parts = cell_part.replace('R', '').split('C')
                    if len(cell_parts) != 2:
                        continue
                    
                    row_idx = int(cell_parts[0]) - 1
                    col_idx = int(cell_parts[1]) - 1
                    
                    if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                        cell = table.rows[row_idx].cells[col_idx]
                        if cell.paragraphs:
                            replace_text_keep_format_docx(cell.paragraphs[0], translated_value)
        
        except (ValueError, IndexError, AttributeError) as e:
            # Bỏ qua các key không hợp lệ
            continue
    
    return doc

@app.route('/login', methods=['GET', 'POST'])
def login():
    """Trang đăng nhập"""
    if request.method == 'POST':
        password = request.form.get('password', '')
        correct_password = get_password()
        
        if password == correct_password:
            session.permanent = True
            session['logged_in'] = True
            session['session_id'] = create_session_id()
            
            # Cleanup old sessions khi đăng nhập
            cleanup_old_sessions()
            
            return redirect(url_for('index'))
        else:
            return render_template('login.html', error='Mật khẩu không đúng!')
    
    return render_template('login.html')

@app.route('/logout')
def logout():
    """Đăng xuất"""
    # Xóa folder của session hiện tại
    if 'session_id' in session:
        session_folder = os.path.join(app.config['UPLOAD_FOLDER'], session['session_id'])
        if os.path.exists(session_folder):
            shutil.rmtree(session_folder)
    
    session.clear()
    return redirect(url_for('login'))

@app.route('/')
@login_required
def index():
    """
    Trang chủ hiển thị dashboard với 2 chức năng Extract và Inject
    """
    # Cleanup old sessions mỗi khi load trang
    cleanup_old_sessions()
    return render_template('index.html')

@app.route('/api/languages', methods=['GET'])
@login_required
def get_languages():
    """
    Trả về danh sách ngôn ngữ đích từ file languages.json
    """
    languages_file = os.path.join(os.path.dirname(__file__), 'languages.json')
    try:
        with open(languages_file, 'r', encoding='utf-8') as f:
            languages = json.load(f)
        return jsonify(languages)
    except FileNotFoundError:
        # Fallback nếu file không tồn tại
        return jsonify([
            {"code": "ja", "name": "tiếng Nhật",  "label": "🇯🇵 Tiếng Nhật (Japanese)"},
            {"code": "en", "name": "tiếng Anh",   "label": "🇺🇸 Tiếng Anh (English)"},
            {"code": "vi", "name": "tiếng Việt",  "label": "🇻🇳 Tiếng Việt (Vietnamese)"},
            {"code": "zh", "name": "tiếng Trung", "label": "🇨🇳 Tiếng Trung (Chinese)"},
            {"code": "ko", "name": "tiếng Hàn",   "label": "🇰🇷 Tiếng Hàn (Korean)"}
        ])
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ==================== API: PROMPT TEMPLATES ====================


@app.route('/api/templates', methods=['GET'])
@login_required
def api_get_templates():
    """Trả về danh sách prompt templates cho ngôn ngữ được chỉ định"""
    lang = request.args.get('lang', 'default')
    if lang == 'img-ocr':
        return jsonify(load_img_ocr_prompt_template())
    return jsonify(load_templates(lang))


@app.route('/api/templates', methods=['POST'])
@login_required
def api_save_templates():
    """Lưu danh sách prompt templates cho ngôn ngữ được chỉ định"""
    lang = request.args.get('lang', 'default')
    new_templates = request.json
    if not isinstance(new_templates, list):
        return jsonify({'error': 'Dữ liệu phải là array'}), 400
    try:
        if lang == 'img-ocr':
            save_img_ocr_prompt_template(new_templates)
            return jsonify({'success': True})
        try:
            with open(TEMPLATES_FILE, 'r', encoding='utf-8') as f:
                all_data = json.load(f)
            if isinstance(all_data, list):
                all_data = {'default': all_data}
        except Exception:
            all_data = {}
        all_data[lang] = new_templates
        with open(TEMPLATES_FILE, 'w', encoding='utf-8') as f:
            json.dump(all_data, f, ensure_ascii=False, indent=2)
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ==================== API: GLOSSARY ====================

@app.route('/api/glossaries', methods=['GET'])
@login_required
def api_list_glossaries():
    """Trả về danh sách tất cả file glossary đã lưu."""
    result = []
    for fname in sorted(os.listdir(GLOSSARY_DIR)):
        if not fname.endswith('.csv'):
            continue
        gid = fname[:-4]
        meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')
        display_name = gid
        if os.path.exists(meta_path):
            with open(meta_path, 'r', encoding='utf-8') as f:
                meta = json.load(f)
                display_name = meta.get('name', gid)
        # Đếm số dòng
        try:
            with open(os.path.join(GLOSSARY_DIR, fname), 'r', encoding='utf-8-sig') as f:
                row_count = sum(1 for r in csv.reader(f) if any(r))
        except Exception:
            row_count = 0
        result.append({'id': gid, 'name': display_name, 'rows': row_count})
    return jsonify(result)


@app.route('/api/glossaries', methods=['POST'])
@login_required
def api_upload_glossary():
    """Upload file CSV mới. Form fields: file (CSV), name (tên hiển thị)."""
    import time
    if 'file' not in request.files:
        return jsonify({'error': 'Không có file'}), 400
    f = request.files['file']
    if not f.filename.endswith('.csv'):
        return jsonify({'error': 'Chỉ chấp nhận file .csv'}), 400
    display_name = request.form.get('name', '').strip() or os.path.splitext(f.filename)[0]
    gid = f'glossary_{int(time.time())}'
    csv_path  = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
    meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')
    f.save(csv_path)
    with open(meta_path, 'w', encoding='utf-8') as mf:
        json.dump({'name': display_name}, mf, ensure_ascii=False)
    # Đếm dòng
    with open(csv_path, 'r', encoding='utf-8-sig') as cf:
        row_count = sum(1 for r in csv.reader(cf) if any(r))
    return jsonify({'success': True, 'id': gid, 'name': display_name, 'rows': row_count})


@app.route('/api/glossaries/<gid>', methods=['GET'])
@login_required
def api_get_glossary(gid):
    """Trả về toàn bộ nội dung glossary dưới dạng list of {src, dst}."""
    csv_path = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
    if not os.path.exists(csv_path):
        return jsonify({'error': 'Không tìm thấy'}), 404
    rows = []
    with open(csv_path, 'r', encoding='utf-8-sig') as f:
        for r in csv.reader(f):
            if len(r) >= 2:
                rows.append({'dst': r[0].strip(), 'src': r[1].strip()})
            elif len(r) == 1 and r[0].strip():
                rows.append({'dst': r[0].strip(), 'src': ''})
    return jsonify(rows)


@app.route('/api/glossaries/<gid>', methods=['PUT'])
@login_required
def api_update_glossary(gid):
    """
    Cập nhật nội dung glossary.
    Body JSON: { "name": "...", "rows": [{"src": "...", "dst": "..."}, ...] }
    """
    csv_path  = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
    meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')
    if not os.path.exists(csv_path):
        return jsonify({'error': 'Không tìm thấy'}), 404
    data = request.get_json()
    if not data:
        return jsonify({'error': 'Body rỗng'}), 400
    # Cập nhật tên
    if 'name' in data:
        with open(meta_path, 'w', encoding='utf-8') as mf:
            json.dump({'name': data['name']}, mf, ensure_ascii=False)
    # Cập nhật rows
    if 'rows' in data:
        with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            for row in data['rows']:
                if row.get('src') or row.get('dst'):
                    w.writerow([row.get('dst', ''), row.get('src', '')])
    return jsonify({'success': True})


@app.route('/api/glossaries/<gid>', methods=['DELETE'])
@login_required
def api_delete_glossary(gid):
    """Xóa glossary và meta file."""
    csv_path  = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
    meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')
    for p in [csv_path, meta_path]:
        try:
            os.remove(p)
        except FileNotFoundError:
            pass
    return jsonify({'success': True})


# ==================== API: EXTRACT COLORS ====================

@app.route('/api/extract-colors', methods=['POST'])
@login_required
def api_extract_colors():
    """
    Trả về danh sách màu chữ duy nhất có trong file xlsx/pptx/docx.
    Luôn bao gồm '000000' để đại diện cho màu đen/auto (mặc định).
    """
    if 'file' not in request.files or not request.files['file'].filename:
        return jsonify({'error': 'Không có file'}), 400

    f = request.files['file']
    if not allowed_file(f.filename):
        return jsonify({'error': 'Chỉ chấp nhận file .xlsx, .pptx hoặc .docx'}), 400

    ext = f.filename.rsplit('.', 1)[1].lower()
    file_bytes = f.read()
    colors = set()

    try:
        if ext == 'xlsx':
            wb = load_workbook(io.BytesIO(file_bytes))
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            c = _get_font_rgb_xlsx(cell)
                            if c:
                                colors.add(c)
            wb.close()

        elif ext == 'pptx':
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    _collect_pptx_shape_colors(shape, colors)

        elif ext == 'docx':
            doc = Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                for run in para.runs:
                    c = _get_font_rgb_docx(run)
                    if c:
                        colors.add(c)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                c = _get_font_rgb_docx(run)
                                if c:
                                    colors.add(c)

    except Exception as e:
        return jsonify({'error': f'Lỗi đọc màu: {str(e)}'}), 500

    # Luôn bao gồm 000000 cho màu đen/auto
    colors.add('000000')

    return jsonify({'colors': sorted(colors)})


@app.route('/extract', methods=['POST'])
@login_required
def extract():
    """
    Trích xuất file → trả về Server-Sent Events (SSE) với progress feedback.
    Events: reading(10%), chunking(40-60%), writing(75-88%), done(100%, result) | error.
    """
    # Phase 1: Xử lý file TRƯỚC khi bắt đầu stream (session phải được set trước khi trả response)
    su_info = session.get('tab1_from_smart_update')
    if 'file' in request.files and request.files['file'].filename:
        file = request.files['file']
        if not allowed_file(file.filename):
            return jsonify({'error': 'Chỉ chấp nhận file .xlsx, .pptx hoặc .docx'}), 400
        original_filename = file.filename
        if '.' not in original_filename:
            return jsonify({'error': 'Tên file phải có đuôi mở rộng'}), 400
        original_ext = original_filename.rsplit('.', 1)[1].lower()
        session_folder = get_session_folder()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        safe_temp_filename = f"temp_{timestamp}.{original_ext}"
        filepath = os.path.join(session_folder, safe_temp_filename)
        file.save(filepath)
    elif su_info and os.path.exists(su_info.get('filepath', '')):
        original_filename = su_info['display_name']
        original_ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else 'xlsx'
        session_folder = get_session_folder()
        filepath = su_info['filepath']
    else:
        return jsonify({'error': 'Không có file được upload'}), 400

    # Thu thập các tham số
    color_filter_raw = request.form.get('color_filter', '')
    color_list = [c.strip() for c in color_filter_raw.split(',') if c.strip()]
    color_filter = _normalize_color_filter(color_list) if color_list else None
    glossary_ids_raw = request.form.get('glossary_ids', '')
    glossary_ids = [g.strip() for g in glossary_ids_raw.split(',') if g.strip()]
    proofread_mode = request.form.get('proofread_mode', '').strip().lower() in ('1', 'true', 'yes', 'on')

    # Tạo session key để inject có thể tìm lại file nguồn (phải set TRƯỚC khi stream)
    session_key = f'sse_extract_{datetime.now().strftime("%Y%m%d_%H%M%S_%f")}'
    session[session_key] = {'filepath': filepath, 'display_name': original_filename}
    # Cũng lưu vào tab1_from_smart_update để tương thích với inject path cũ
    session['tab1_from_smart_update'] = {'filepath': filepath, 'display_name': original_filename}

    resp = Response(
        stream_with_context(stream_extract(
            filepath, original_filename, glossary_ids, session_folder, color_filter, proofread_mode=proofread_mode
        )),
        mimetype='text/event-stream',
    )
    resp.headers['Cache-Control'] = 'no-cache'
    resp.headers['X-Accel-Buffering'] = 'no'
    resp.headers['X-Session-Key'] = session_key
    return resp



# ==================== HELPER: raw extract (no chunking/zip) ====================

def _extract_raw(filepath, original_filename, glossary_ids, session_folder, color_filter=None, selected_sheets=None, proofread_mode=False):
    """
    Trích xuất raw text từ file, trả về dict {key: value}.
    Không tạo ZIP hay chunk - chỉ trích xuất data thô và áp glossary.
    """
    original_ext = original_filename.rsplit('.', 1)[-1].lower() if '.' in original_filename else 'xlsx'

    if original_ext == 'xlsx':
        workbook = load_workbook(filepath)
        extracted_data = {}
        for sheet_name in workbook.sheetnames:
            if selected_sheets and sheet_name not in selected_sheets:
                continue
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str) and not cell.value.startswith('='):
                        if color_filter is None or (_get_font_rgb_xlsx(cell) or '000000') in color_filter:
                            extracted_data[f"{sheet_name}!{cell.coordinate}"] = cell.value
        extracted_data.update(extract_xlsx_shapes(filepath))
        workbook.close()
    elif original_ext == 'pptx':
        extracted_data = extract_text_from_pptx(filepath, color_filter)
    elif original_ext == 'docx':
        extracted_data = extract_text_from_docx(filepath, color_filter)
    else:
        raise ValueError(f'Không hỗ trợ định dạng .{original_ext}')

    if proofread_mode:
        extracted_data = _filter_proofread_extract_data(extracted_data)

    if glossary_ids:
        extracted_data = apply_glossary(extracted_data, glossary_ids)

    return extracted_data


# ==================== HELPER: core extract logic ====================

def _run_extract(filepath, original_filename, glossary_ids, session_folder, color_filter=None, selected_sheets=None, proofread_mode=False):
    """
    Chạy toàn bộ logic extract từ cột filepath.
    Trả về dict cho jsonify (cùng format như route /extract).
    Ném Exception nếu có lỗi.
    color_filter: set HEX strings hoặc None (không lọc)
    selected_sheets: list tên sheet muốn extract, hoặc None (tất cả)
    """
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    original_ext = original_filename.rsplit('.', 1)[-1].lower() if '.' in original_filename else 'xlsx'

    if original_ext == 'xlsx':
        workbook = load_workbook(filepath)
        extracted_data = {}
        for sheet_name in workbook.sheetnames:
            if selected_sheets and sheet_name not in selected_sheets:
                continue
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str) and not cell.value.startswith('='):
                        if color_filter is None or (_get_font_rgb_xlsx(cell) or '000000') in color_filter:
                            extracted_data[f"{sheet_name}!{cell.coordinate}"] = cell.value
        # TODO: color filter for xlsx shapes not yet implemented
        extracted_data.update(extract_xlsx_shapes(filepath))
        workbook.close()
    elif original_ext == 'pptx':
        extracted_data = extract_text_from_pptx(filepath, color_filter)
    elif original_ext == 'docx':
        extracted_data = extract_text_from_docx(filepath, color_filter)
    else:
        raise ValueError(f'Không hỗ trợ định dạng .{original_ext}')

    if proofread_mode:
        extracted_data = _filter_proofread_extract_data(extracted_data)

    if glossary_ids:
        extracted_data = apply_glossary(extracted_data, glossary_ids)

    CHUNK_SIZE = 400
    data_items  = list(extracted_data.items())
    total_items = len(data_items)
    num_files   = max(1, (total_items + CHUNK_SIZE - 1) // CHUNK_SIZE)

    base_filename = os.path.splitext(original_filename)[0] or f'file_{timestamp}'
    safe_base     = f'extracted_{timestamp}'
    folder_name   = f'{base_filename}_json_to_translate'
    temp_dir      = os.path.join(session_folder, f'{safe_base}_temp_{timestamp}')
    os.makedirs(temp_dir, exist_ok=True)

    json_files = []
    json_display_names = []
    for i in range(num_files):
        chunk_data  = dict(data_items[i*CHUNK_SIZE:(i+1)*CHUNK_SIZE])
        disp_name   = f'{base_filename}_part{i+1:02d}_of_{num_files:02d}.json'
        safe_name   = f'{safe_base}_part{i+1:02d}.json'
        jpath       = os.path.join(temp_dir, safe_name)
        with open(jpath, 'w', encoding='utf-8') as jf:
            json.dump(chunk_data, jf, ensure_ascii=False, indent=2)
        json_files.append(jpath)
        json_display_names.append(disp_name)

    files_data = []
    for idx, jpath in enumerate(json_files):
        with open(jpath, 'r', encoding='utf-8') as f:
            files_data.append({'name': json_display_names[idx], 'content': f.read()})

    zip_display = f'{base_filename}_json_to_translate.zip'
    safe_zip    = f'{safe_base}_json_{timestamp}.zip'
    zip_path    = os.path.join(session_folder, safe_zip)
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_STORED) as zf:
        for idx, jpath in enumerate(json_files):
            zf.write(jpath, os.path.join(folder_name, json_display_names[idx]))

    session['extract_zip'] = {
        'path': zip_path,
        'display_name': zip_display,
        'input_path': filepath,
        'json_files': json_files,
        'temp_dir': temp_dir,
    }

    dedup_files, dedup_mapping, dedup_stats = build_dedup_data(extracted_data, CHUNK_SIZE)
    dedup_map_path = os.path.join(session_folder, 'dedup_mapping.json')
    with open(dedup_map_path, 'w', encoding='utf-8') as f:
        json.dump(dedup_mapping, f, ensure_ascii=False, indent=2)

    return {
        'success': True,
        'total_files': num_files,
        'total_items': total_items,
        'files': files_data,
        'zip_display_name': zip_display,
        'dedup_files': dedup_files,
        'dedup_stats': dedup_stats,
    }


@app.route('/load-google-sheet', methods=['POST'])
@login_required
def load_google_sheet():
    """
    Nhận link Google Sheet, tải file xlsx về, lưu vào session folder.
    Trả về: { session_key, display_name, sheets }
    """
    data = request.get_json() or {}
    url  = data.get('url', '').strip()
    if not url:
        return jsonify({'error': 'Thiếu link Google Sheet'}), 400

    try:
        spreadsheet_id, _ = parse_google_sheet_url(url)
    except ValueError as e:
        return jsonify({'error': str(e)}), 400

    export_url = build_sheet_export_url(spreadsheet_id)
    try:
        resp = _requests.get(export_url, timeout=30, allow_redirects=True)
        ct   = resp.headers.get('Content-Type', '')
        if resp.status_code != 200 or (
            'spreadsheet' not in ct and 'officedocument' not in ct and len(resp.content) < 1000
        ):
            return jsonify({'error': (
                'Không thể tải file. Hãy kiểm tra file đã được share công khai '
                '(Anyone with link → View).'
            )}), 400
    except Exception as e:
        return jsonify({'error': f'Lỗi kết nối: {e}'}), 500

    session_folder = get_session_folder()
    filename = f'gsheet_{spreadsheet_id[:12]}.xlsx'
    filepath = os.path.join(session_folder, filename)
    with open(filepath, 'wb') as f:
        f.write(resp.content)

    try:
        wb = load_workbook(filepath, read_only=True)
        sheet_names = wb.sheetnames
        wb.close()
    except Exception:
        sheet_names = []

    session_key = f'gsheet_{uuid.uuid4().hex[:8]}'
    session[session_key] = {
        'filepath':     filepath,
        'display_name': f'GoogleSheet_{spreadsheet_id[:12]}.xlsx',
        'sheets':       sheet_names,
    }

    return jsonify({
        'session_key':  session_key,
        'display_name': f'GoogleSheet_{spreadsheet_id[:12]}.xlsx',
        'sheets':       sheet_names,
    })


@app.route('/extract-from-sheet', methods=['POST'])
@login_required
def extract_from_sheet():
    """
    Extract nội dung từ Google Sheet đã tải về (lưu trong session).
    Nhận: { session_key, selected_sheets, glossary_ids }
    """
    data            = request.get_json() or {}
    session_key     = data.get('session_key')
    glossary_ids    = data.get('glossary_ids', [])
    selected_sheets = data.get('selected_sheets') or None  # None = tất cả

    if not session_key or session_key not in session:
        return jsonify({'error': 'Phiên làm việc hết hạn. Vui lòng tải lại Google Sheet.'}), 400

    info     = session[session_key]
    filepath = info['filepath']
    if not os.path.exists(filepath):
        return jsonify({'error': 'File tạm không còn tồn tại. Vui lòng tải lại Google Sheet.'}), 400

    try:
        session_folder = get_session_folder()
        result = _run_extract(filepath, info['display_name'], glossary_ids, session_folder, selected_sheets=selected_sheets)
        return jsonify(result)
    except Exception as e:
        return jsonify({'error': f'Lỗi khi xử lý sheet: {str(e)}'}), 500


@app.route('/download-zip', methods=['GET'])
@login_required
def download_zip():
    """
    Serve file ZIP đã được tạo từ /extract.
    Xóa tất cả file tạm sau khi gửi xong.
    """
    zip_info = session.get('extract_zip')
    if not zip_info:
        # Fallback: kiểm tra file trạng thái trên filesystem (SSE extract path)
        try:
            sf = get_session_folder()
            state_path = os.path.join(sf, 'extract_state.json')
            if os.path.exists(state_path):
                with open(state_path, 'r', encoding='utf-8') as _f:
                    zip_info = json.load(_f)
                os.remove(state_path)
        except Exception:
            pass
    if not zip_info:
        return jsonify({'error': 'Không tìm thấy file ZIP. Vui lòng trích xuất lại.'}), 404
    
    zip_filepath = zip_info.get('path')
    zip_display_name = zip_info.get('display_name', 'download.zip')
    
    if not zip_filepath or not os.path.exists(zip_filepath):
        return jsonify({'error': 'File ZIP không còn tồn tại. Vui lòng trích xuất lại.'}), 404
    
    # Xóa thông tin ZIP trong session
    session.pop('extract_zip', None)
    
    # Trả về file ZIP
    response = send_file(zip_filepath, mimetype='application/zip')
    response = set_download_headers(response, zip_display_name, 'download.zip')
    
    # Xóa tất cả file tạm sau khi gửi
    input_path = zip_info.get('input_path')
    json_files = zip_info.get('json_files', [])
    temp_dir = zip_info.get('temp_dir')
    
    @response.call_on_close
    def cleanup():
        import time
        import gc
        gc.collect()
        time.sleep(0.1)
        
        try:
            if zip_filepath and os.path.exists(zip_filepath):
                os.remove(zip_filepath)
        except Exception as e:
            print(f"Warning: Không thể xóa ZIP: {e}")
        
        try:
            if input_path and os.path.exists(input_path):
                os.remove(input_path)
        except Exception as e:
            print(f"Warning: Không thể xóa input file: {e}")
        
        for jf in json_files:
            try:
                if os.path.exists(jf):
                    os.remove(jf)
            except Exception as e:
                print(f"Warning: Không thể xóa JSON file: {e}")
        
        try:
            if temp_dir and os.path.exists(temp_dir):
                os.rmdir(temp_dir)
        except Exception as e:
            print(f"Warning: Không thể xóa temp dir: {e}")
    
    return response

@app.route('/inject', methods=['POST'])
@login_required
def inject():
    """
    Chức năng 2: Nạp dữ liệu từ file JSON đã dịch vào file Excel, PPTX hoặc DOCX gốc
    Giữ nguyên định dạng, màu sắc của file gốc
    Hỗ trợ nhiều file JSON riêng lẻ hoặc file ZIP chứa nhiều file JSON
    """
    # Kiểm tra file upload hoặc fallback từ Smart Update hoặc Google Sheet
    su_info_inject = session.get('tab1_from_smart_update')
    sheet_session_key = request.form.get('sheet_session_key')
    
    if 'excel_file' in request.files and request.files['excel_file'].filename:
        excel_file = request.files['excel_file']
        use_session_file_inject = False
        if not allowed_file(excel_file.filename):
            return jsonify({'error': 'File phải có định dạng .xlsx, .pptx hoặc .docx'}), 400
    elif sheet_session_key and sheet_session_key in session:
        # FIX: Lấy file từ Google Sheet session
        sheet_info = session[sheet_session_key]
        if os.path.exists(sheet_info.get('filepath', '')):
            excel_file = None
            use_session_file_inject = True
            # Đặt filepath từ sheet session thay vì smart update
            su_info_inject = sheet_info
        else:
            return jsonify({'error': 'File Sheet đã hết hạn hoặc không tồn tại. Vui lòng tải lại Google Sheet.'}), 400
    elif su_info_inject and os.path.exists(su_info_inject.get('filepath', '')):
        excel_file = None
        use_session_file_inject = True
    else:
        return jsonify({'error': 'Cần upload file Excel, PPTX hoặc DOCX'}), 400

    # Lấy pasted JSON data nếu có
    pasted_json_data = request.form.get('pasted_json_data', None)

    # Kiểm tra xem có file JSON được upload hoặc có pasted JSON không
    json_files = request.files.getlist('json_files') if 'json_files' in request.files else []

    # Kiểm tra xem có ít nhất một nguồn JSON
    has_json_files = len(json_files) > 0 and any(f.filename != '' for f in json_files)
    has_pasted_json = pasted_json_data is not None and pasted_json_data.strip() != ''

    if not has_json_files and not has_pasted_json:
        return jsonify({'error': 'Cần upload ít nhất 1 file JSON/ZIP hoặc paste JSON'}), 400

    try:
        # Lấy session folder
        session_folder = get_session_folder()
        
        if not use_session_file_inject:
            # Lưu tên file gốc (giữ nguyên tiếng Nhật, ký tự đặc biệt)
            original_excel_filename = excel_file.filename

            # Lấy extension từ tên file gốc
            if '.' in original_excel_filename:
                original_ext = original_excel_filename.rsplit('.', 1)[1].lower()
            else:
                return jsonify({'error': 'Tên file phải có đuôi mở rộng (.xlsx hoặc .pptx)'}), 400

            # Tạo tên file tạm an toàn hoàn toàn từ timestamp
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            safe_temp_filename = f"temp_{timestamp}.{original_ext}"
            excel_filepath = os.path.join(session_folder, safe_temp_filename)
            excel_file.save(excel_filepath)
        else:
            original_excel_filename = su_info_inject['display_name']
            original_ext = original_excel_filename.rsplit('.', 1)[1].lower() if '.' in original_excel_filename else 'xlsx'
            excel_filepath = su_info_inject['filepath']
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            session.pop('tab1_from_smart_update', None)

        # Đọc và gộp dữ liệu JSON từ tất cả các file
        json_data = {}
        temp_files = []
        
        for json_file in json_files:
            if json_file.filename == '':
                continue
                
            json_filename = json_file.filename.lower()
            is_zip = json_filename.endswith('.zip')
            is_json = json_filename.endswith('.json')
            
            if not (is_zip or is_json):
                continue
            
            if is_zip:
                # Xử lý file ZIP
                temp_zip_filename = f"temp_{timestamp}_{secure_filename(json_file.filename)}"
                zip_filepath = os.path.join(session_folder, temp_zip_filename)
                json_file.save(zip_filepath)
                temp_files.append(zip_filepath)
                
                # Giải nén và đọc tất cả các file JSON
                with zipfile.ZipFile(zip_filepath, 'r') as zipf:
                    for file_info in zipf.namelist():
                        if file_info.lower().endswith('.json'):
                            with zipf.open(file_info) as f:
                                try:
                                    content = f.read().decode('utf-8')
                                    chunk_data = json.loads(content)
                                    json_data.update(chunk_data)
                                except UnicodeDecodeError:
                                    # Thử với encoding khác
                                    f.seek(0)
                                    content = f.read().decode('utf-8-sig')
                                    chunk_data = json.loads(content)
                                    json_data.update(chunk_data)
            else:
                # Xử lý file JSON đơn lẻ
                try:
                    json_content = json_file.stream.read().decode('utf-8')
                    chunk_data = json.loads(json_content)
                    json_data.update(chunk_data)
                except UnicodeDecodeError:
                    # Thử với encoding khác nếu UTF-8 thất bại
                    json_file.stream.seek(0)
                    json_content = json_file.stream.read().decode('utf-8-sig')
                    chunk_data = json.loads(json_content)
                    json_data.update(chunk_data)
                except json.JSONDecodeError as e:
                    return jsonify({'error': f'File JSON "{json_file.filename}" không hợp lệ: {str(e)}'}), 400
        
        # Xử lý pasted JSON data
        if pasted_json_data:
            try:
                pasted_data_list = json.loads(pasted_json_data)
                
                # pasted_data_list là danh sách các JSON objects
                if isinstance(pasted_data_list, list):
                    for idx, pasted_obj in enumerate(pasted_data_list):
                        if isinstance(pasted_obj, dict):
                            json_data.update(pasted_obj)
                        else:
                            return jsonify({'error': f'Pasted JSON #{idx + 1} không phải là object'}), 400
                else:
                    return jsonify({'error': 'Pasted JSON data phải là danh sách các objects'}), 400
                    
            except json.JSONDecodeError as e:
                return jsonify({'error': f'Pasted JSON không hợp lệ: {str(e)}'}), 400

        # Mở rộng dedup keys nếu có (khi user dịch từ dedup JSON)
        if any(k.startswith('dedup_') for k in json_data):
            json_data = expand_dedup_data(json_data, session_folder)

        # Xác định loại file và nạp dữ liệu (dùng tên file gốc)
        file_ext = original_excel_filename.rsplit('.', 1)[1].lower()

        # Tạo backup trước khi inject (người dùng có thể tải về nếu cần)
        try:
            bk_ts = datetime.now().strftime('%Y%m%d_%H%M%S')
            backup_path = os.path.join(session_folder, f'backup_{bk_ts}.{file_ext}')
            shutil.copy2(excel_filepath, backup_path)
            session['last_inject_backup'] = {
                'path': backup_path,
                'display_name': f'backup_{original_excel_filename}',
                'created': datetime.now().isoformat(),
            }
        except Exception:
            pass  # Backup thất bại không chặn inject

        if file_ext == 'xlsx':
            # Tạo tên file output
            base_filename = os.path.splitext(original_excel_filename)[0]  # Tên gốc với tiếng Nhật
            
            output_display_name = f"{base_filename}_translated.xlsx"  # Tên hiển thị
            safe_output_filename = f"output_{timestamp}.xlsx"  # Tên file trong filesystem
            output_filepath = os.path.join(session_folder, safe_output_filename)

            # ZIP-level patch: nạp cả cell + shape trực tiếp trên package gốc
            inject_xlsx_shapes(excel_filepath, output_filepath, json_data)
            
            output_mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        
        elif file_ext == 'pptx':
            # Nạp text vào PPTX
            prs = inject_text_to_pptx(excel_filepath, json_data)
            
            # Tạo tên file output
            base_filename = os.path.splitext(original_excel_filename)[0]  # Tên gốc với tiếng Nhật
            
            output_display_name = f"{base_filename}_translated.pptx"  # Tên hiển thị
            safe_output_filename = f"output_{timestamp}.pptx"  # Tên file trong filesystem
            output_filepath = os.path.join(session_folder, safe_output_filename)
            
            # Lưu file PPTX đã được nạp dữ liệu
            prs.save(output_filepath)
            del prs  # Giải phóng memory
            
            output_mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        elif file_ext == 'docx':
            # Nạp text vào DOCX
            doc = inject_text_to_docx(excel_filepath, json_data)
            
            # Tạo tên file output
            base_filename = os.path.splitext(original_excel_filename)[0]  # Tên gốc với tiếng Nhật
            
            output_display_name = f"{base_filename}_translated.docx"  # Tên hiển thị
            safe_output_filename = f"output_{timestamp}.docx"  # Tên file trong filesystem
            output_filepath = os.path.join(session_folder, safe_output_filename)
            
            # Lưu file DOCX đã được nạp dữ liệu
            doc.save(output_filepath)
            del doc  # Giải phóng memory
            
            output_mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        
        # Trả về file đã được nạp dữ liệu (dùng tên hiển thị)
        response = send_file(
            output_filepath,
            mimetype=output_mimetype
        )
        
        default_ascii_name = 'download.docx' if file_ext == 'docx' else ('download.pptx' if file_ext == 'pptx' else 'download.xlsx')
        response = set_download_headers(response, output_display_name, default_ascii_name)
        
        # FIX: Pop sheet_session_key khỏi session sau khi inject thành công
        if sheet_session_key:
            session.pop(sheet_session_key, None)
        
        # Xóa tất cả file tạm sau khi gửi response
        @response.call_on_close
        def cleanup():
            import time
            import gc
            
            # Force garbage collection để giải phóng file handles
            gc.collect()
            time.sleep(0.1)  # Delay nhỏ để đảm bảo file được giải phóng
            
            # Xóa file output
            try:
                if os.path.exists(output_filepath):
                    os.remove(output_filepath)
            except Exception as e:
                print(f"Warning: Không thể xóa output file: {e}")
            
            # Xóa file Excel/PPTX/DOCX tạm
            try:
                if os.path.exists(excel_filepath):
                    os.remove(excel_filepath)
            except Exception as e:
                print(f"Warning: Không thể xóa file tạm: {e}")
            
            # Xóa tất cả file ZIP tạm
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception as e:
                    print(f"Warning: Không thể xóa file ZIP tạm: {e}")
        
        return response
        
    except Exception as e:
        # Xử lý lỗi
        return jsonify({'error': f'Lỗi khi xử lý file: {str(e)}'}), 500


# ==================== PROOF-MAP HELPERS ====================

def proof_map_xlsx(source_path, output_path, json_data, hex_color):
    """Write proof-read output for xlsx: original text + corrected text in hex_color."""
    from openpyxl.styles import Alignment as _Alignment
    from openpyxl.cell.rich_text import CellRichText as _CellRichText, TextBlock as _TextBlock
    from openpyxl.cell.text import InlineFont as _InlineFont
    wb = load_workbook(source_path)
    corr_argb = 'FF' + hex_color.lstrip('#').upper()

    for key, corrected_val in json_data.items():
        try:
            if '!' not in key:
                continue
            sheet_name, coord = key.split('!', 1)
            if sheet_name not in wb.sheetnames:
                continue
            ws = wb[sheet_name]
            cell = ws[coord]
            orig_val = cell.value
            if orig_val is None:
                orig_val = ''
            if str(orig_val).strip() == str(corrected_val).strip():
                continue

            # Preserve original font color (ARGB format)
            try:
                fc = cell.font.color if cell.font else None
                orig_argb = (fc.rgb if fc and fc.type == 'rgb' and fc.rgb else 'FF000000')
            except Exception:
                orig_argb = 'FF000000'

            # Use openpyxl rich text API so Excel renders text runs correctly.
            try:
                rich_text = _CellRichText()
                rich_text.append(_TextBlock(_InlineFont(color=orig_argb), str(orig_val)))
                rich_text.append(_TextBlock(_InlineFont(color=corr_argb), ' ' + str(corrected_val)))
                cell.value = rich_text
            except Exception:
                # Safe fallback: keep content even if rich-text object is not supported.
                cell.value = f"{orig_val} {corrected_val}"

            cell.alignment = _Alignment(wrap_text=True)
        except Exception:
            continue

    wb.save(output_path)


def _proof_navigate_pptx_shape(shape, shape_indices):
    """Navigate nested grouped shapes; returns final shape or None."""
    if not shape_indices:
        return shape
    if hasattr(shape, 'shapes'):
        idx = shape_indices[0]
        if idx <= len(shape.shapes):
            return _proof_navigate_pptx_shape(shape.shapes[idx - 1], shape_indices[1:])
    return None


def proof_map_pptx(source_path, output_path, json_data, hex_color):
    """Write proof-read output for pptx: appends colored correction paragraph."""
    from pptx.dml.color import RGBColor as _PptxRGB

    def _pick_ref_run(text_frame):
        """Pick a reference run to inherit font properties (name/size)."""
        fallback_run = None
        for para in text_frame.paragraphs:
            for run in para.runs:
                if fallback_run is None:
                    fallback_run = run
                if run.font and (run.font.name or run.font.size):
                    return run
        return fallback_run

    r_c = int(hex_color[1:3], 16)
    g_c = int(hex_color[3:5], 16)
    b_c = int(hex_color[5:7], 16)
    corr_rgb = _PptxRGB(r_c, g_c, b_c)
    prs = Presentation(source_path)

    for key, corrected_text in json_data.items():
        try:
            if '!' not in key:
                continue
            parts = key.split('!')
            if len(parts) < 2:
                continue
            slide_part = parts[0]
            if not slide_part.startswith('Slide'):
                continue
            slide_idx = int(slide_part.replace('Slide', '')) - 1
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]

            shape_part = parts[1]
            if not shape_part.startswith('Shape'):
                continue
            shape_str = shape_part.replace('Shape', '')
            shape_indices = [int(i) for i in shape_str.split('_')]
            first_idx = shape_indices[0] - 1
            if first_idx >= len(slide.shapes):
                continue
            shape = slide.shapes[first_idx]
            shape = _proof_navigate_pptx_shape(shape, shape_indices[1:])
            if shape is None:
                continue

            is_table_cell = len(parts) == 3 and parts[2].startswith('Table_R')
            if is_table_cell:
                table_part = parts[2].replace('Table_R', '').split('C')
                row_idx = int(table_part[0]) - 1
                col_idx = int(table_part[1]) - 1
                if not (hasattr(shape, 'has_table') and shape.has_table):
                    continue
                table = shape.table
                if row_idx >= len(table.rows) or col_idx >= len(table.rows[row_idx].cells):
                    continue
                text_frame = table.rows[row_idx].cells[col_idx].text_frame
            else:
                if not hasattr(shape, 'text_frame') or not shape.text_frame:
                    continue
                text_frame = shape.text_frame

            orig_text = text_frame.text
            if orig_text.strip() == corrected_text.strip():
                continue

            # Append correction run to last paragraph (same line, no new paragraph)
            last_para = text_frame.paragraphs[-1]
            new_run = last_para.add_run()
            new_run.text = ' ' + corrected_text
            new_run.font.color.rgb = corr_rgb

            # Keep font family and size consistent with original text.
            try:
                ref_run = _pick_ref_run(text_frame)
                if ref_run and ref_run.font:
                    if ref_run.font.name:
                        new_run.font.name = ref_run.font.name
                    if ref_run.font.size:
                        new_run.font.size = ref_run.font.size
            except Exception:
                pass
        except Exception:
            continue

    prs.save(output_path)


def proof_map_docx(source_path, output_path, json_data, hex_color):
    """Write proof-read output for docx: inserts colored correction paragraph after original."""
    from docx.shared import RGBColor as _DocxRGB
    r_c = int(hex_color[1:3], 16)
    g_c = int(hex_color[3:5], 16)
    b_c = int(hex_color[5:7], 16)
    corr_rgb = _DocxRGB(r_c, g_c, b_c)
    doc = Document(source_path)

    def _insert_correction(paragraph, corrected_text):
        # Append correction as a new run in the same paragraph and inherit font size/name.
        new_run = paragraph.add_run(' ' + corrected_text)
        new_run.font.color.rgb = corr_rgb

        ref_run = None
        for run in paragraph.runs:
            if run is new_run:
                continue
            if ref_run is None:
                ref_run = run
            if run.font and (run.font.name or run.font.size):
                ref_run = run
                break

        if ref_run and ref_run.font:
            if ref_run.font.name:
                new_run.font.name = ref_run.font.name
            if ref_run.font.size:
                new_run.font.size = ref_run.font.size

    for key, corrected_text in json_data.items():
        try:
            # 1. Plain paragraph: "ParagraphX"
            if key.startswith('Paragraph') and '!' not in key:
                para_num = int(key.replace('Paragraph', ''))
                para_idx = 0
                for para in doc.paragraphs:
                    if para.text.strip():
                        para_idx += 1
                        if para_idx == para_num:
                            if para.text.strip() != corrected_text.strip():
                                _insert_correction(para, corrected_text)
                            break

            # 2. Table cell: "TableX!RyCz"
            elif key.startswith('Table') and '!' in key and not key.startswith('Header_') and not key.startswith('Footer_'):
                parts = key.split('!')
                if len(parts) != 2:
                    continue
                table_idx = int(parts[0].replace('Table', '')) - 1
                if table_idx >= len(doc.tables):
                    continue
                table = doc.tables[table_idx]
                cell_str = parts[1]
                if not cell_str.startswith('R'):
                    continue
                rc = cell_str.replace('R', '').split('C')
                row_idx = int(rc[0]) - 1
                col_idx = int(rc[1]) - 1
                if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                    cell = table.rows[row_idx].cells[col_idx]
                    if cell.paragraphs:
                        para = cell.paragraphs[0]
                        if para.text.strip() != corrected_text.strip():
                            _insert_correction(para, corrected_text)

            # 3. Header: "Header_SectionX!ParagraphY" or "Header_SectionX!TableY!RzCw"
            elif key.startswith('Header_Section'):
                parts = key.split('!')
                if len(parts) < 2:
                    continue
                sec_idx = int(parts[0].replace('Header_Section', '')) - 1
                if sec_idx >= len(doc.sections):
                    continue
                header = doc.sections[sec_idx].header
                if parts[1].startswith('Paragraph'):
                    pn = int(parts[1].replace('Paragraph', ''))
                    pi = 0
                    for para in header.paragraphs:
                        if para.text.strip():
                            pi += 1
                            if pi == pn:
                                if para.text.strip() != corrected_text.strip():
                                    _insert_correction(para, corrected_text)
                                break
                elif parts[1].startswith('Table') and len(parts) == 3:
                    ti = int(parts[1].replace('Table', '')) - 1
                    if ti >= len(header.tables):
                        continue
                    rc = parts[2].replace('R', '').split('C')
                    ri, ci = int(rc[0]) - 1, int(rc[1]) - 1
                    if ri < len(header.tables[ti].rows) and ci < len(header.tables[ti].rows[ri].cells):
                        para = header.tables[ti].rows[ri].cells[ci].paragraphs[0]
                        if para.text.strip() != corrected_text.strip():
                            _insert_correction(para, corrected_text)

            # 4. Footer: "Footer_SectionX!ParagraphY" or "Footer_SectionX!TableY!RzCw"
            elif key.startswith('Footer_Section'):
                parts = key.split('!')
                if len(parts) < 2:
                    continue
                sec_idx = int(parts[0].replace('Footer_Section', '')) - 1
                if sec_idx >= len(doc.sections):
                    continue
                footer = doc.sections[sec_idx].footer
                if parts[1].startswith('Paragraph'):
                    pn = int(parts[1].replace('Paragraph', ''))
                    pi = 0
                    for para in footer.paragraphs:
                        if para.text.strip():
                            pi += 1
                            if pi == pn:
                                if para.text.strip() != corrected_text.strip():
                                    _insert_correction(para, corrected_text)
                                break
                elif parts[1].startswith('Table') and len(parts) == 3:
                    ti = int(parts[1].replace('Table', '')) - 1
                    if ti >= len(footer.tables):
                        continue
                    rc = parts[2].replace('R', '').split('C')
                    ri, ci = int(rc[0]) - 1, int(rc[1]) - 1
                    if ri < len(footer.tables[ti].rows) and ci < len(footer.tables[ti].rows[ri].cells):
                        para = footer.tables[ti].rows[ri].cells[ci].paragraphs[0]
                        if para.text.strip() != corrected_text.strip():
                            _insert_correction(para, corrected_text)
        except Exception:
            continue

    doc.save(output_path)


# ==================== PROOF-MAP ROUTES ====================

@app.route('/proof-map', methods=['POST'])
@login_required
def proof_map():
    """
    Map kết quả kiểm tra ngữ pháp vào file gốc.
    Keys không thay đổi → giữ nguyên.
    Keys có sửa → thêm dòng mới với màu correction_color bên dưới.
    """
    # 1. Get source file (same as /inject: check session then upload)
    su_info = session.get('tab1_from_smart_update')
    if 'excel_file' in request.files and request.files['excel_file'].filename:
        excel_file = request.files['excel_file']
        use_session = False
        if not allowed_file(excel_file.filename):
            return jsonify({'error': 'Chỉ chấp nhận .xlsx, .pptx, .docx'}), 400
    elif su_info and os.path.exists(su_info.get('filepath', '')):
        use_session = True
        excel_file = None
    else:
        return jsonify({'error': 'Cần upload file gốc'}), 400

    # 2. Get correction color
    hex_color = request.form.get('correction_color', '#FF0000').strip()
    if not re.match(r'^#[0-9A-Fa-f]{6}$', hex_color):
        hex_color = '#FF0000'

    # 3. Get pasted JSON
    pasted_json_data = request.form.get('pasted_json_data', '').strip()
    if not pasted_json_data:
        return jsonify({'error': 'Cần paste JSON kết quả từ AI'}), 400
    try:
        json_data = json.loads(pasted_json_data)
    except json.JSONDecodeError as e:
        return jsonify({'error': f'JSON không hợp lệ: {e}'}), 400

    # 4. Expand dedup keys if present
    session_folder = get_session_folder()
    if any(k.startswith('dedup_') for k in json_data):
        json_data = expand_dedup_data(json_data, session_folder)

    # 5. Save source file to temp if uploaded
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    try:
        if not use_session:
            original_filename = excel_file.filename
            ext = original_filename.rsplit('.', 1)[1].lower()
            src_path = os.path.join(session_folder, f'proof_src_{timestamp}.{ext}')
            excel_file.save(src_path)
        else:
            original_filename = su_info['display_name']
            ext = original_filename.rsplit('.', 1)[1].lower()
            src_path = su_info['filepath']

        # 6. Output path
        base_name = os.path.splitext(original_filename)[0]
        out_display = f'{base_name}_proofread.{ext}'
        out_safe    = f'proof_out_{timestamp}.{ext}'
        out_path    = os.path.join(session_folder, out_safe)

        # 7. Dispatch to helper
        if ext == 'xlsx':
            proof_map_xlsx(src_path, out_path, json_data, hex_color)
            mimetype = ('application/vnd.openxmlformats-officedocument'
                        '.spreadsheetml.sheet')
        elif ext == 'pptx':
            proof_map_pptx(src_path, out_path, json_data, hex_color)
            mimetype = ('application/vnd.openxmlformats-officedocument'
                        '.presentationml.presentation')
        elif ext == 'docx':
            proof_map_docx(src_path, out_path, json_data, hex_color)
            mimetype = ('application/vnd.openxmlformats-officedocument'
                        '.wordprocessingml.document')
        else:
            return jsonify({'error': 'Định dạng không hỗ trợ'}), 400

        # 8. Store in session for download
        token = uuid.uuid4().hex[:12]
        session[f'proof_dl_{token}'] = {
            'path': out_path,
            'display_name': out_display,
            'mimetype': mimetype,
        }
        return jsonify({
            'success': True,
            'download_token': token,
            'display_name': out_display,
        })

    except Exception as e:
        app.logger.error(f'proof_map error: {e}', exc_info=True)
        return jsonify({'error': f'Lỗi xử lý: {str(e)}'}), 500


@app.route('/download-proof/<token>', methods=['GET'])
@login_required
def download_proof(token):
    """Serve file kết quả kiểm tra ngữ pháp."""
    key = f'proof_dl_{token}'
    if key not in session:
        return jsonify({'error': 'Token không hợp lệ hoặc đã hết hạn'}), 404
    info = session.pop(key)
    path = info['path']
    if not os.path.exists(path):
        return jsonify({'error': 'File không còn tồn tại'}), 404
    response = send_file(path, mimetype=info['mimetype'])
    response = set_download_headers(
        response, info['display_name'], 'proofread_output.' + path.rsplit('.', 1)[-1])
    return response


@app.route('/download-inject-backup', methods=['GET'])
@login_required
def download_inject_backup():
    """Tải về file gốc đã được backup trước khi inject."""
    backup_info = session.get('last_inject_backup')
    if not backup_info or not os.path.exists(backup_info.get('path', '')):
        return jsonify({'error': 'Không có backup nào. Hãy thực hiện inject trước.'}), 404
    ext = backup_info['path'].rsplit('.', 1)[-1].lower()
    mimetypes_map = {
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    }
    mimetype = mimetypes_map.get(ext, 'application/octet-stream')
    response = send_file(backup_info['path'], mimetype=mimetype)
    response = set_download_headers(response, backup_info['display_name'], f'backup.{ext}')
    return response


@app.route('/clear-uploads', methods=['POST'])
@login_required
def clear_uploads():
    """
    Xóa tất cả file trong thư mục session hiện tại
    """
    try:
        session_folder = get_session_folder()
        
        # Kiểm tra xem thư mục có tồn tại không
        if not os.path.exists(session_folder):
            return jsonify({'success': True, 'message': 'Không có file nào để xóa'}), 200
        
        # Đếm số file đã xóa
        deleted_count = 0
        
        # Duyệt qua tất cả file trong session folder
        for item in os.listdir(session_folder):
            item_path = os.path.join(session_folder, item)
            
            try:
                if os.path.isfile(item_path):
                    # Xóa file
                    os.remove(item_path)
                    deleted_count += 1
                elif os.path.isdir(item_path):
                    # Xóa thư mục con và tất cả nội dung bên trong
                    shutil.rmtree(item_path)
                    deleted_count += 1
            except Exception as e:
                print(f"Không thể xóa {item_path}: {str(e)}")
        
        return jsonify({
            'success': True,
            'message': f'Đã xóa thành công {deleted_count} file trong phiên của bạn',
            'deleted_count': deleted_count
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Lỗi khi xóa file: {str(e)}'}), 500

# ==================== SMART UPDATE ROUTES ====================

@app.route('/smart-update', methods=['POST'])
@login_required
def smart_update_route():
    """
    Smart Update: So sánh và kế thừa bản dịch Excel từ version cũ.
    Input (form-data):
        file_vn_old / file_vn10  : VN_1.0 (chưa dịch, phiên bản cũ)
        file_vn_new / file_vn11  : VN_1.1 (chưa dịch, phiên bản mới)
        file_jp_old / file_jp10  : JP_1.0 (đã dịch, phiên bản cũ)
        new_colors (optional)    : danh sách RGB6 cách nhau dấu phẩy, VD "38761D,00B050"
    Output JSON: stats + link tải file
    """
    # Hỗ trợ cả 2 bộ tên trường (cũ và mới)
    def _get_file(new_name, old_name):
        f = request.files.get(new_name) or request.files.get(old_name)
        return f if (f and f.filename) else None

    file_vn10 = _get_file('file_vn_old', 'file_vn10')
    file_vn11 = _get_file('file_vn_new', 'file_vn11')
    file_jp10 = _get_file('file_jp_old', 'file_jp10')

    missing = [n for n, f in [('VN cũ', file_vn10), ('VN mới', file_vn11), ('JP cũ', file_jp10)] if not f]
    if missing:
        return jsonify({'error': f'Thiếu file: {", ".join(missing)}'}), 400

    for f in [file_vn10, file_vn11, file_jp10]:
        if not f.filename.lower().endswith('.xlsx'):
            return jsonify({'error': 'Smart Update chỉ hỗ trợ file .xlsx'}), 400

    # Parse màu xanh tùy chỉnh từ form (nếu có)
    raw_colors = request.form.get('new_colors', '').strip()
    custom_new_colors = None
    if raw_colors:
        custom_new_colors = set(
            c.strip().upper().lstrip('#').lstrip('FF')[-6:].zfill(6)
            for c in raw_colors.split(',') if c.strip()
        ) | DEFAULT_NEW_COLORS  # merge với mặc định

    try:
        session_folder = get_session_folder()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        path_vn10 = os.path.join(session_folder, f'su_vn10_{timestamp}.xlsx')
        path_vn11 = os.path.join(session_folder, f'su_vn11_{timestamp}.xlsx')
        path_jp10 = os.path.join(session_folder, f'su_jp10_{timestamp}.xlsx')
        file_vn10.save(path_vn10)
        file_vn11.save(path_vn11)
        file_jp10.save(path_jp10)

        result_wb, to_translate, stats = smart_update_excel(
            path_vn10, path_vn11, path_jp10,
            new_colors=custom_new_colors
        )

        # Lưu workbook kết quả JP_1.1
        original_name = os.path.splitext(file_vn11.filename)[0]
        result_display_name = f"{original_name}_JP_1_1.xlsx"
        safe_result_path = os.path.join(session_folder, f'su_result_{timestamp}.xlsx')
        result_wb.save(safe_result_path)
        result_wb.close()

        # Chia các ô cần dịch thành JSON chunks (≤400/file)
        CHUNK_SIZE = 400
        items = list(to_translate.items())
        num_chunks = max(1, (len(items) + CHUNK_SIZE - 1) // CHUNK_SIZE)

        temp_json_dir = os.path.join(session_folder, f'su_json_{timestamp}')
        os.makedirs(temp_json_dir, exist_ok=True)

        json_display_names = []
        json_paths = []
        files_data = []

        for i in range(num_chunks):
            chunk = dict(items[i * CHUNK_SIZE:(i + 1) * CHUNK_SIZE])
            jname_display = f"{original_name}_to_translate_part{i+1:02d}_of_{num_chunks:02d}.json"
            jname_safe = f"su_json_part{i+1:02d}_{timestamp}.json"
            jpath = os.path.join(temp_json_dir, jname_safe)
            with open(jpath, 'w', encoding='utf-8') as jf:
                json.dump(chunk, jf, ensure_ascii=False, indent=2)
            json_display_names.append(jname_display)
            json_paths.append(jpath)
            with open(jpath, 'r', encoding='utf-8') as jf:
                files_data.append({'name': jname_display, 'content': jf.read()})

        # Tạo ZIP chứa tất cả JSON cần dịch
        zip_display_name = f"{original_name}_to_translate.zip"
        zip_path = os.path.join(session_folder, f'su_json_{timestamp}.zip')
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_STORED) as zf:
            for idx, jpath_item in enumerate(json_paths):
                zf.write(jpath_item, json_display_names[idx])

        session['smart_update'] = {
            'result_path': safe_result_path,
            'result_display_name': result_display_name,
            'result_filename': os.path.basename(safe_result_path),
            'path_vn11': path_vn11,
            'zip_path': zip_path,
            'zip_display_name': zip_display_name,
            'temp_files': [path_vn10, path_vn11, path_jp10] + json_paths,
            'temp_dir': temp_json_dir,
        }

        return jsonify({
            'success': True,
            'stats': stats,
            'result_display_name': result_display_name,
            'zip_display_name': zip_display_name,
            'files': files_data,
        })

    except Exception as e:
        import traceback
        return jsonify({'error': f'Lỗi xử lý: {str(e)}', 'detail': traceback.format_exc()}), 500


@app.route('/download-smart-excel', methods=['GET'])
@login_required
def download_smart_excel():
    """Tải file Excel JP 1.1 hỗn hợp từ Smart Update"""
    info = session.get('smart_update')
    if not info or not os.path.exists(info.get('result_path', '')):
        return jsonify({'error': 'Không tìm thấy file. Vui lòng chạy Smart Update lại.'}), 404
    response = send_file(
        info['result_path'],
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response = set_download_headers(response, info['result_display_name'], 'result_mixed.xlsx')
    return response


@app.route('/download-smart-zip', methods=['GET'])
@login_required
def download_smart_zip():
    """Tải ZIP các file JSON cần dịch từ Smart Update và dọn dẹp file tạm"""
    info = session.get('smart_update')
    if not info or not os.path.exists(info.get('zip_path', '')):
        return jsonify({'error': 'Không tìm thấy file ZIP. Vui lòng chạy Smart Update lại.'}), 404

    zip_path = info['zip_path']
    zip_display_name = info['zip_display_name']
    session.pop('smart_update', None)

    response = send_file(zip_path, mimetype='application/zip')
    response = set_download_headers(response, zip_display_name, 'new_strings.zip')

    temp_files = info.get('temp_files', [])
    temp_dir = info.get('temp_dir')
    result_path = info.get('result_path')

    @response.call_on_close
    def cleanup():
        import time
        import gc
        gc.collect()
        time.sleep(0.1)
        for p in [zip_path, result_path] + temp_files:
            try:
                if p and os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
        try:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
        except Exception:
            pass

    return response


@app.route('/smart-update/use-as-input', methods=['POST'])
@login_required
def smart_update_use_as_input():
    """
    Chuyển file kết quả Smart Update (JP_1.1) thành file input cho Tab Dịch tài liệu.
    Frontend gọi route này, sau đó chuyển sang Tab 1 — file đã sẵn sàng để Extract/Inject.
    """
    info = session.get('smart_update')
    if not info or not os.path.exists(info.get('result_path', '')):
        return jsonify({'error': 'Không tìm thấy file Smart Update. Vui lòng chạy lại.'}), 404

    result_path  = info['result_path']
    display_name = info['result_display_name']

    # Ghi vào session của Tab 1 — giống như user vừa upload file này
    session['tab1_from_smart_update'] = {
        'filepath':     result_path,
        'display_name': display_name,
    }

    return jsonify({
        'success':      True,
        'display_name': display_name,
    })


# Đảm bảo thư mục uploads tồn tại
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])


# ==================== IMAGE TRANSLATION ====================
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp'}

def allowed_image(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_IMAGE_EXTENSIONS


@app.route('/img-translate/upload', methods=['POST'])
@login_required
def img_translate_upload():
    """Upload ảnh vào session folder, trả về filename."""
    if 'image' not in request.files:
        return jsonify({'error': 'Không tìm thấy file ảnh'}), 400
    file = request.files['image']
    if file.filename == '':
        return jsonify({'error': 'Chưa chọn file'}), 400
    if not allowed_image(file.filename):
        return jsonify({'error': 'Định dạng không hỗ trợ. Chỉ nhận: PNG, JPG, JPEG, GIF, WEBP, BMP'}), 400

    session_folder = get_session_folder()
    filename = secure_filename(file.filename)
    ts = datetime.now().strftime('%H%M%S')
    name, ext = os.path.splitext(filename)
    safe_filename = f"img_{ts}_{name[:30]}{ext}"
    filepath = os.path.join(session_folder, safe_filename)
    file.save(filepath)
    return jsonify({'success': True, 'filename': safe_filename}), 200


@app.route('/img-translate/image/<filename>', methods=['GET'])
@login_required
def img_translate_serve(filename):
    """Serve ảnh đã upload từ session folder."""
    session_folder = get_session_folder()
    safe = secure_filename(filename)
    filepath = os.path.join(session_folder, safe)
    if not os.path.exists(filepath):
        return jsonify({'error': 'File không tồn tại'}), 404
    return send_file(filepath)


# NEW: OCR route — calls OCR.space API for exact bounding boxes instead of AI coordinate guessing
@app.route('/img-translate/ocr', methods=['POST'])
@login_required
def img_translate_ocr():
    """Call OCR.space API on the uploaded image. Returns text blocks with pixel bounding boxes."""
    data = request.get_json() or {}
    filename = data.get('filename', '').strip()
    # NEW: OCR.space language code (jpn, chs, cht, kor, eng, ara, fre, ger, por, rus, spa, tur)
    ocr_lang = data.get('ocr_lang', 'jpn')

    if not filename:
        return jsonify({'error': 'Thiếu tên file'}), 400

    session_folder = get_session_folder()
    safe = secure_filename(filename)
    filepath = os.path.join(session_folder, safe)
    if not os.path.exists(filepath):
        return jsonify({'error': 'File không tồn tại'}), 404

    # NEW: image natural dimensions passed from frontend for pixel coord context
    img_w = int(data.get('img_w', 0))
    img_h = int(data.get('img_h', 0))

    # NEW: OCR.space free public endpoint — 500 req/day per IP with helloworld key
    OCR_API_URL = 'https://api.ocr.space/parse/image'

    try:
        with open(filepath, 'rb') as f:
            # NEW: isOverlayRequired=true is MANDATORY to receive bounding box coordinates
            resp = _requests.post(
                OCR_API_URL,
                files={'file': (safe, f, 'image/png')},
                data={
                    'apikey': 'helloworld',
                    'language': ocr_lang,
                    'isOverlayRequired': 'true',
                    'detectOrientation': 'true',
                    # FIX B: scale=false so bbox coords stay in original image space
                    'scale': 'false',
                    'OCREngine': '2',
                    'isCreateSearchablePdf': 'false',
                    'isSearchablePdfHideTextLayer': 'false',
                },
                timeout=30
            )
        result = resp.json()
    except Exception as e:
        return jsonify({'error': f'Lỗi gọi OCR.space: {str(e)}'}), 500

    # DEBUG: log top-level response structure
    # app.logger.debug(f'[OCR DEBUG] Response keys: {list(result.keys())}')
    # app.logger.debug(f'[OCR DEBUG] IsErroredOnProcessing: {result.get("IsErroredOnProcessing")}')
    # app.logger.debug(f'[OCR DEBUG] ParsedResults count: {len(result.get("ParsedResults", []))}')

    if result.get('IsErroredOnProcessing'):
        err_msg = result.get('ErrorMessage', ['Lỗi không xác định'])
        if isinstance(err_msg, list):
            err_msg = ' '.join(err_msg)
        return jsonify({'error': f'OCR.space lỗi: {err_msg}'}), 500

    # NEW: parse OCR.space Lines→Words structure into our block format
    blocks = []
    parsed = result.get('ParsedResults', [])
    # Keep raw_iw/raw_ih in outer scope so the return debug dict can reference them
    raw_iw = raw_ih = None
    ocr_img_w = ocr_img_h = 0
    for page in parsed:
        overlay = page.get('TextOverlay', {})
        # FIX: read dimensions reported by OCR.space for this page
        raw_iw = overlay.get('ImageWidth')
        raw_ih = overlay.get('ImageHeight')
        ocr_img_w = int(raw_iw) if raw_iw and int(raw_iw) > 0 else 0
        ocr_img_h = int(raw_ih) if raw_ih and int(raw_ih) > 0 else 0

        # FIX STEP 5: when OCR.space doesn't report dimensions, infer from Engine 2 normalization
        # Engine 2 rescales the long side to max 1024 px; bbox coords are in that scaled space
        if ocr_img_w == 0 or ocr_img_h == 0:
            if img_w > 0 and img_h > 0:
                max_side = max(img_w, img_h)
                if max_side > 1024:
                    scale_factor = 1024.0 / max_side
                    ocr_img_w = round(img_w * scale_factor)
                    ocr_img_h = round(img_h * scale_factor)
                else:
                    # Already within Engine 2 limit — no rescaling applied
                    ocr_img_w = img_w
                    ocr_img_h = img_h
            else:
                ocr_img_w = img_w or 1024
                ocr_img_h = img_h or 1024

        # DEBUG: log dimension resolution for each page
        # app.logger.debug(
        #     f'[OCR DEBUG] TextOverlay raw ImageWidth={raw_iw}, ImageHeight={raw_ih}'
        # )
        # app.logger.debug(
        #     f'[OCR DEBUG] Frontend sent img_w={img_w}, img_h={img_h}'
        # )
        # app.logger.debug(
        #     f'[OCR DEBUG] Resolved reference: ocr_img_w={ocr_img_w}, ocr_img_h={ocr_img_h}'
        # )
        # app.logger.debug(
        #     f'[OCR DEBUG] Lines count: {len(overlay.get("Lines", []))}'
        # )
        lines   = overlay.get('Lines', [])
        for line_idx, line in enumerate(lines):
            words = line.get('Words', [])
            if not words:
                continue
            line_text = ' '.join(w.get('WordText', '') for w in words).strip()
            if not line_text:
                continue

            # Build union bounding box from individual word boxes
            lefts   = [w['Left']             for w in words if 'Left'   in w]
            tops    = [w['Top']              for w in words if 'Top'    in w]
            rights  = [w['Left'] + w['Width']  for w in words if 'Left'  in w and 'Width'  in w]
            bottoms = [w['Top']  + w['Height'] for w in words if 'Top'   in w and 'Height' in w]
            if not lefts:
                continue

            x0 = min(lefts);  y0 = min(tops)
            x1 = max(rights); y1 = max(bottoms)

            # FIX 1: OCR.space sometimes emits MinTop=0/MinLeft=0 for a line even when the
            # actual text is not at the image origin.  When both x0 and y0 are 0, try to use
            # the first word's coordinates as a more reliable anchor.
            line_min_top  = line.get('MinTop',  y0)
            line_min_left = line.get('MinLeft', x0)
            if line_min_top == 0 and line_min_left == 0:
                first_word = next(
                    (w for w in words if w.get('Left', -1) > 0 or w.get('Top', -1) > 0),
                    None
                )
                if first_word:
                    # Re-anchor: shift x0/y0 to first reliable word; keep x1/y1 from union
                    fw_left = first_word.get('Left', x0)
                    fw_top  = first_word.get('Top',  y0)
                    if fw_left > 0 or fw_top > 0:
                        x0 = min(fw_left, x0) if x0 > 0 else fw_left
                        y0 = min(fw_top,  y0) if y0 > 0 else fw_top
                        # app.logger.debug(
                        #     f'[OCR DEBUG] Line {line_idx} "{line_text[:20]}" had MinTop/MinLeft=0; '
                        #     f're-anchored to first word at ({fw_left},{fw_top})'
                        # )

            # FIX 2: validation — skip blocks that still resolve to (0,0) when the image is
            # clearly larger, as these are almost certainly corrupt OCR entries that would
            # overlap with any real corner text.
            has_valid_size = (x1 - x0) > 2 and (y1 - y0) > 2
            at_true_corner = x0 == 0 and y0 == 0 and ocr_img_w > 0 and ocr_img_h > 0
            if at_true_corner and ocr_img_w > 50 and ocr_img_h > 50:
                # Accept only if the block is a plausible corner word (small width/height)
                plausible_corner = (x1 - x0) < ocr_img_w * 0.3 and (y1 - y0) < ocr_img_h * 0.15
                if not plausible_corner:
                    # app.logger.warning(
                    #     f'[OCR DEBUG] Skipping line {line_idx} "{line_text[:30]}" — '
                    #     f'(0,0) origin with suspicious size {x1-x0}×{y1-y0} px'
                    # )
                    continue

            if not has_valid_size:
                # app.logger.debug(f'[OCR DEBUG] Skipping line {line_idx} — zero-size bbox')
                continue

            # FIX 3/4: use average word height (not full line height) for font sizing.
            # Line bbox often includes inter-line spacing; word heights reflect actual glyphs.
            word_heights = sorted([w['Height'] for w in words if 'Height' in w and w['Height'] > 0])
            if word_heights:
                avg_word_h = sum(word_heights) / len(word_heights)
                # Prefer average over median for font size — median can be skewed by ascenders
                word_h = round(avg_word_h)
            else:
                word_h = y1 - y0

            # DEBUG: log each block's raw pixel coords and computed percentages
            # app.logger.debug(
            #     f"[OCR DEBUG] Block '{line_text[:30]}': "
            #     f"pixel x={x0},y={y0},w={x1-x0},h={y1-y0} | word_h={word_h} | "
            #     f"pct top={round(y0/ocr_img_h*100,2)}, left={round(x0/ocr_img_w*100,2)}, "
            #     f"w={round((x1-x0)/ocr_img_w*100,2)}, h={round((y1-y0)/ocr_img_h*100,2)}"
            # )
            blocks.append({
                'text':   line_text,
                'x':      x0,
                'y':      y0,
                'w':      x1 - x0,
                'h':      y1 - y0,
                'word_h': word_h,   # average word glyph height in OCR pixel space
                # FIX: OCR.space returns coords in the NATURAL image coordinate space,
                # not the normalized (1024px) space. The frontend divides pixel coords by
                # img_w/img_h to get %, so we must use the original natural dimensions here.
                'img_w':  img_w if img_w > 0 else ocr_img_w,
                'img_h':  img_h if img_h > 0 else ocr_img_h,
            })

    if not blocks:
        return jsonify({'error': 'Không nhận diện được văn bản. Thử chọn ngôn ngữ OCR khác.'}), 400

    # FIX: natural_img_w/h are what the frontend uses as the % denominator
    natural_img_w = img_w if img_w > 0 else ocr_img_w
    natural_img_h = img_h if img_h > 0 else ocr_img_h
    app.logger.debug(
        f'[OCR DEBUG] Natural img: {img_w}\u00d7{img_h}, '
        f'OCR normalized: {ocr_img_w}\u00d7{ocr_img_h}, '
        f'Using as block ref: img_w={natural_img_w}'
    )
    # DEBUG: return dimension resolution info alongside blocks for browser console inspection
    return jsonify({
        'blocks': blocks,
        'count':  len(blocks),
        'debug': {
            'ocr_normalized_w':        ocr_img_w,
            'ocr_normalized_h':        ocr_img_h,
            'natural_img_w':           natural_img_w,
            'natural_img_h':           natural_img_h,
            'frontend_img_w':          img_w,
            'frontend_img_h':          img_h,
            'imagewidth_from_response': raw_iw,
            'imageheight_from_response': raw_ih,
            # FIX: compare against natural dims, not normalized dims
            'dimensions_match':        (natural_img_w == img_w and natural_img_h == img_h),
        }
    }), 200


@app.route('/img-translate/prompt', methods=['POST'])
@login_required
def img_translate_prompt():
    """Sinh Instruction Prompt yêu cầu AI phân tích ảnh và trả JSON overlay (fallback khi OCR.space không được)."""
    data = request.get_json() or {}
    target_lang = data.get('target_lang', 'tiếng Nhật')
    source_lang = data.get('source_lang', '').strip()

    # NEW: removed img_w/img_h — AI internal resize makes pixel->% conversion wrong
    source_note = f" (ngôn ngữ gốc trong ảnh: {source_lang})" if source_lang else ""

    # NEW: removed bg_color, text_color, font_size_pct — frontend computes from canvas + geometry
    prompt = f"""Bạn là chuyên gia OCR và dịch thuật chuyên nghiệp. Tôi sẽ gửi cho bạn một bức ảnh{source_note}.

NHIỆM VỤ:
1. Nhận diện (OCR) TẤT CẢ các vùng có văn bản trong ảnh.
2. Dịch toàn bộ sang {target_lang} một cách tự nhiên, chính xác.
3. Với mỗi vùng văn bản, xác định:

   top_pct    : tọa độ mép TRÊN của text box, % chiều CAO ảnh bạn nhìn thấy (0–100)
   left_pct   : tọa độ mép TRÁI của text box, % chiều RỘNG ảnh bạn nhìn thấy (0–100)
   width_pct  : chiều RỘNG text box, % chiều rộng ảnh bạn nhìn thấy
   height_pct : chiều CAO text box, % chiều cao ảnh bạn nhìn thấy

⚠️ YÊU CẦU BẮT BUỘC:
- Tọa độ % tính theo kích thước ảnh BẠN ĐANG XỬ LÝ
- Tọa độ phải bao phủ CHÍNH XÁC vùng chứa văn bản, sai số không quá 1%
- Trả về JSON THUẦN TÚY — KHÔNG giải thích, KHÔNG bọc markdown code block ```

FORMAT JSON TRẢ VỀ (giữ đúng cấu trúc này):
{{
  "text_blocks": [
    {{
      "original": "Văn bản gốc trong ảnh",
      "translated": "Bản dịch sang {target_lang}",
      "top_pct": 10.5,
      "left_pct": 5.2,
      "width_pct": 30.0,
      "height_pct": 5.0
    }}
  ]
}}"""

    return jsonify({'prompt': prompt}), 200


# ==================== TERMINOLOGY EXTRACTOR ====================

TERMINOLOGY_PROMPTS_FILE = 'terminology_prompts.json'


def get_default_terminology_prompts() -> list:
    """Trả về danh sách prompt mặc định cho việc trích xuất thuật ngữ."""
    return [
        {
            "id": "terminology_default",
            "name": "Trích xuất thuật ngữ chuyên ngành",
            "content": (
                "Bạn là chuyên gia ngôn ngữ chuyên ngành. Tôi sẽ cung cấp một danh sách "
                "các cặp văn bản song ngữ dạng JSON (mỗi cặp gồm \"src\": văn bản gốc, \"dst\": văn bản đã dịch).\n\n"
                "Nhiệm vụ:\n"
                "1. Phân tích toàn bộ các cặp văn bản\n"
                "2. Xác định các từ/cụm từ CHUYÊN NGÀNH xuất hiện nhất quán trong bản dịch\n"
                "3. Chỉ lấy thuật ngữ thực sự chuyên ngành — loại bỏ từ thông thường, giới từ, động từ phổ thông\n"
                "4. Loại bỏ trùng lặp (giữ lại 1 cặp duy nhất cho mỗi thuật ngữ)\n\n"
                "Trả về KẾT QUẢ là JSON object thuần túy theo format sau:\n"
                "{\n"
                "  \"src_text_1\": \"dst_text_1\",\n"
                "  \"src_text_2\": \"dst_text_2\"\n"
                "}\n\n"
                "Trong đó:\n"
                "- KEY = thuật ngữ ngôn ngữ GỐC (src)\n"
                "- VALUE = thuật ngữ ngôn ngữ ĐÍCH (dst)\n\n"
                "⚠️ QUAN TRỌNG:\n"
                "- Chỉ trả về JSON object thuần túy, không markdown, không giải thích\n"
                "- Không bọc trong ```json``` hay bất kỳ code block nào\n"
                "- Đảm bảo JSON hợp lệ, không có trailing comma\n"
                "- Nếu không tìm thấy thuật ngữ chuyên ngành, trả về {}"
            )
        }
    ]


def load_terminology_prompts() -> list:
    try:
        with open(TERMINOLOGY_PROMPTS_FILE, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return data if isinstance(data, list) and data else get_default_terminology_prompts()
    except Exception:
        return get_default_terminology_prompts()


def save_terminology_prompts(prompts: list) -> None:
    with open(TERMINOLOGY_PROMPTS_FILE, 'w', encoding='utf-8') as f:
        json.dump(prompts, f, ensure_ascii=False, indent=2)


def align_bilingual_texts(src_dict: dict, dst_dict: dict) -> list:
    """
    Ghép hai dict extract theo key chung.
    Chỉ lấy các key có mặt trong CẢ HAI và cả hai value đều không rỗng.
    Loại bỏ các cặp có src == dst (không thực sự được dịch).
    Trả về list of {"src": str, "dst": str}.
    """
    import unicodedata
    pairs = []
    common_keys = set(src_dict.keys()) & set(dst_dict.keys())
    for key in sorted(common_keys):
        src_text = unicodedata.normalize('NFC', str(src_dict[key]).strip())
        dst_text = unicodedata.normalize('NFC', str(dst_dict[key]).strip())
        if src_text and dst_text and src_text != dst_text:
            pairs.append({"src": src_text, "dst": dst_text})
    return pairs


def extract_text_from_file(filepath: str, ext: str) -> dict:
    """
    Wrapper tái sử dụng logic extract hiện có.
    ext: 'xlsx' | 'pptx' | 'docx'
    Trả về dict {key: text}.
    """
    if ext == 'xlsx':
        wb = load_workbook(filepath, data_only=True)
        result = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str) and not cell.value.startswith('='):
                        result[f"{sheet_name}!{cell.coordinate}"] = cell.value
        result.update(extract_xlsx_shapes(filepath))
        wb.close()
        return result
    elif ext == 'pptx':
        return extract_text_from_pptx(filepath)
    elif ext == 'docx':
        return extract_text_from_docx(filepath)
    return {}


@app.route('/api/terminology/align', methods=['POST'])
@login_required
def api_terminology_align():
    """
    Nhận file_src + file_dst, extract cả hai, ghép cặp song ngữ theo key.
    """
    if 'file_src' not in request.files or 'file_dst' not in request.files:
        return jsonify({'error': 'Cần upload cả file gốc (file_src) và file đã dịch (file_dst)'}), 400

    f_src = request.files['file_src']
    f_dst = request.files['file_dst']

    if not f_src.filename or not f_dst.filename:
        return jsonify({'error': 'Tên file không hợp lệ'}), 400

    if not allowed_file(f_src.filename) or not allowed_file(f_dst.filename):
        return jsonify({'error': 'Chỉ chấp nhận file .xlsx, .pptx hoặc .docx'}), 400

    ext_src = f_src.filename.rsplit('.', 1)[1].lower()
    ext_dst = f_dst.filename.rsplit('.', 1)[1].lower()
    if ext_src != ext_dst:
        return jsonify({'error': f'Hai file phải cùng định dạng (file gốc: .{ext_src}, file dịch: .{ext_dst})'}), 400

    session_folder = get_session_folder()
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    path_src = os.path.join(session_folder, f'term_src_{ts}.{ext_src}')
    path_dst = os.path.join(session_folder, f'term_dst_{ts}.{ext_dst}')

    try:
        f_src.save(path_src)
        f_dst.save(path_dst)

        src_dict = extract_text_from_file(path_src, ext_src)
        dst_dict = extract_text_from_file(path_dst, ext_dst)

        pairs = align_bilingual_texts(src_dict, dst_dict)

        # Lưu pairs vào file tạm (tránh giới hạn ~4KB của Flask session cookie)
        session_folder = get_session_folder()
        pairs_cache_path = os.path.join(session_folder, 'terminology_pairs_cache.json')
        try:
            with open(pairs_cache_path, 'w', encoding='utf-8') as _f:
                json.dump(pairs, _f, ensure_ascii=False)
            session['terminology_pairs_cache'] = pairs_cache_path
        except Exception:
            pass  # non-critical

        pairs_json_str = json.dumps(pairs, ensure_ascii=False, indent=2)

        return jsonify({
            'success': True,
            'pairs': pairs[:50],         # trả về tối đa 50 cho preview
            'total': len(src_dict),
            'matched': len(pairs),
            'pairs_json': pairs_json_str,
        })
    except Exception as e:
        return jsonify({'error': f'Lỗi khi xử lý file: {str(e)}'}), 500
    finally:
        for p in [path_src, path_dst]:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass


@app.route('/api/terminology/prompts', methods=['GET'])
@login_required
def api_get_terminology_prompts():
    return jsonify(load_terminology_prompts())


@app.route('/api/terminology/prompts', methods=['POST'])
@login_required
def api_save_terminology_prompts_route():
    data = request.get_json()
    if not isinstance(data, list):
        return jsonify({'error': 'Phải là array'}), 400
    try:
        save_terminology_prompts(data)
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/terminology/import-csv', methods=['POST'])
@login_required
def api_terminology_import_csv():
    """
    Parse CSV text từ AI output và lưu vào glossary (mới hoặc merge).
    Body JSON: { csv_text, target_glossary_id?, new_glossary_name? }
    """
    import time as _time
    data = request.get_json()
    if not data:
        return jsonify({'error': 'Body rỗng'}), 400

    csv_text = (data.get('csv_text') or '').strip()
    if not csv_text:
        return jsonify({'error': 'CSV rỗng'}), 400

    target_gid = (data.get('target_glossary_id') or '').strip()
    new_name = (data.get('new_glossary_name') or '').strip()

    # Parse CSV — thử dấu phẩy trước, fallback tab
    rows_raw = []
    for line in csv_text.splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            parsed = list(csv.reader([line]))[0]
        except Exception:
            parsed = [line]

        if len(parsed) >= 2:
            rows_raw.append((parsed[0].strip(), parsed[1].strip()))
        elif len(parsed) == 1 and '\t' in parsed[0]:
            parts = parsed[0].split('\t', 1)
            rows_raw.append((parts[0].strip(), parts[1].strip()))
        # else: bỏ qua dòng không đủ cột

    valid_rows = [(dst, src) for dst, src in rows_raw if dst and src]
    if not valid_rows:
        return jsonify({'error': 'Không tìm thấy dữ liệu hợp lệ trong CSV'}), 400

    added = 0
    skipped = 0

    if target_gid:
        # Merge vào glossary có sẵn
        csv_path = os.path.join(GLOSSARY_DIR, f'{target_gid}.csv')
        if not os.path.exists(csv_path):
            return jsonify({'error': f'Không tìm thấy glossary: {target_gid}'}), 404

        # Đọc existing rows để dedup
        existing_set = set()
        existing_rows = []
        try:
            with open(csv_path, 'r', encoding='utf-8-sig') as f:
                for row in csv.reader(f):
                    if len(row) >= 2:
                        existing_rows.append(row)
                        existing_set.add((row[1].strip().lower(), row[0].strip().lower()))
        except Exception:
            pass

        new_rows = []
        for dst, src in valid_rows:
            key = (src.lower(), dst.lower())
            if key in existing_set:
                skipped += 1
            else:
                new_rows.append([dst, src])
                existing_set.add(key)
                added += 1

        with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            for r in existing_rows:
                w.writerow(r)
            for r in new_rows:
                w.writerow(r)

        # Lấy tên từ meta
        meta_path = os.path.join(GLOSSARY_DIR, f'{target_gid}.meta.json')
        gid_name = target_gid
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                gid_name = json.load(f).get('name', target_gid)
        except Exception:
            pass

        return jsonify({
            'success': True,
            'gid': target_gid,
            'name': gid_name,
            'added': added,
            'skipped_duplicate': skipped,
            'total_rows': len(existing_rows) + added,
        })

    else:
        # Tạo glossary mới
        if not new_name:
            new_name = f'Thuật ngữ {datetime.now().strftime("%Y-%m-%d %H:%M")}'

        gid = f'glossary_{int(_time.time())}'
        csv_path = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
        meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')

        seen = set()
        with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            for dst, src in valid_rows:
                key = (src.lower(), dst.lower())
                if key in seen:
                    skipped += 1
                else:
                    w.writerow([dst, src])
                    seen.add(key)
                    added += 1

        with open(meta_path, 'w', encoding='utf-8') as f:
            json.dump({'name': new_name}, f, ensure_ascii=False)

        return jsonify({
            'success': True,
            'gid': gid,
            'name': new_name,
            'added': added,
            'skipped_duplicate': skipped,
            'total_rows': added,
        })


@app.route('/api/terminology/import-json', methods=['POST'])
@login_required
def api_terminology_import_json():
    """
    Parse JSON object {src: dst} từ AI output và lưu vào glossary (mới hoặc merge).
    Body JSON: { json_text: str, target_glossary_id?: str, new_glossary_name?: str }
    """
    import time as _time
    data = request.get_json()
    if not data:
        return jsonify({'error': 'Body rỗng'}), 400

    json_text = (data.get('json_text') or '').strip()
    if not json_text:
        return jsonify({'error': 'JSON rỗng'}), 400

    target_gid = (data.get('target_glossary_id') or '').strip()
    new_name = (data.get('new_glossary_name') or '').strip()

    # Strip markdown code fence nếu AI bọc kết quả
    clean_text = json_text.strip()
    if clean_text.startswith('```'):
        lines = clean_text.splitlines()
        lines = [l for l in lines if not l.strip().startswith('```')]
        clean_text = '\n'.join(lines).strip()

    try:
        parsed = json.loads(clean_text)
    except json.JSONDecodeError as e:
        return jsonify({'error': f'JSON không hợp lệ: {str(e)}'}), 400

    if not isinstance(parsed, dict):
        return jsonify({'error': 'Dữ liệu phải là JSON object {src: dst}'}), 400

    # Chuyển thành list (cột A=dst, cột B=src — chuẩn glossary CSV)
    valid_rows = []
    for src_text, dst_text in parsed.items():
        src = str(src_text).strip()
        dst = str(dst_text).strip()
        if src and dst:
            valid_rows.append((dst, src))

    if not valid_rows:
        return jsonify({'error': 'Không tìm thấy cặp thuật ngữ hợp lệ trong JSON'}), 400

    added = 0
    skipped = 0

    if target_gid:
        csv_path = os.path.join(GLOSSARY_DIR, f'{target_gid}.csv')
        if not os.path.exists(csv_path):
            return jsonify({'error': f'Không tìm thấy glossary: {target_gid}'}), 404

        existing_set = set()
        existing_rows = []
        try:
            with open(csv_path, 'r', encoding='utf-8-sig') as f:
                for row in csv.reader(f):
                    if len(row) >= 2:
                        existing_rows.append(row)
                        existing_set.add((row[1].strip().lower(), row[0].strip().lower()))
        except Exception:
            pass

        new_rows = []
        for dst, src in valid_rows:
            key = (src.lower(), dst.lower())
            if key in existing_set:
                skipped += 1
            else:
                new_rows.append([dst, src])
                existing_set.add(key)
                added += 1

        with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            for r in existing_rows:
                w.writerow(r)
            for r in new_rows:
                w.writerow(r)

        meta_path = os.path.join(GLOSSARY_DIR, f'{target_gid}.meta.json')
        gid_name = target_gid
        try:
            with open(meta_path, 'r', encoding='utf-8') as f:
                gid_name = json.load(f).get('name', target_gid)
        except Exception:
            pass

        return jsonify({
            'success': True, 'gid': target_gid, 'name': gid_name,
            'added': added, 'skipped_duplicate': skipped,
            'total_rows': len(existing_rows) + added,
        })

    else:
        if not new_name:
            new_name = f'Thuật ngữ {datetime.now().strftime("%Y-%m-%d %H:%M")}'
        gid = f'glossary_{int(_time.time())}'
        csv_path = os.path.join(GLOSSARY_DIR, f'{gid}.csv')
        meta_path = os.path.join(GLOSSARY_DIR, f'{gid}.meta.json')

        seen = set()
        with open(csv_path, 'w', encoding='utf-8-sig', newline='') as f:
            w = csv.writer(f)
            for dst, src in valid_rows:
                key = (src.lower(), dst.lower())
                if key not in seen:
                    w.writerow([dst, src])
                    seen.add(key)
                    added += 1
                else:
                    skipped += 1

        with open(meta_path, 'w', encoding='utf-8') as f:
            json.dump({'name': new_name}, f, ensure_ascii=False)

        return jsonify({
            'success': True, 'gid': gid, 'name': new_name,
            'added': added, 'skipped_duplicate': skipped,
            'total_rows': added,
        })


# ==================== BATCH PROCESSING ====================

@app.route('/batch-extract', methods=['POST'])
@login_required
def batch_extract():
    """
    Trích xuất nhiều file cùng lúc với cross-file dedup.
    Gộp tất cả values unique từ mọi file → 1 JSON duy nhất để dịch.
    Input (form-data): files[] + glossary_ids (comma-sep)
    Output JSON: { batch_id, files:[{name,items}], total_items, dedup_stats, dedup_chunks, zip_display_name }
    """
    uploaded_files = request.files.getlist('files')
    valid_files = [f for f in uploaded_files if f.filename]
    if not valid_files:
        return jsonify({'error': 'Không có file được upload'}), 400
    for f in valid_files:
        if not allowed_file(f.filename):
            return jsonify({'error': f'File "{f.filename}" không hợp lệ. Chỉ chấp nhận .xlsx, .pptx, .docx'}), 400

    glossary_ids_raw = request.form.get('glossary_ids', '')
    glossary_ids = [g.strip() for g in glossary_ids_raw.split(',') if g.strip()]

    color_filter_raw = request.form.get('color_filter', '')
    color_filter_list = [c.strip() for c in color_filter_raw.split(',') if c.strip()] if color_filter_raw else None

    session_folder = get_session_folder()
    batch_id = uuid.uuid4().hex[:10]

    batch_session_files = []
    file_extracted_list = []  # [{original_filename, extracted_data}]

    for idx, f in enumerate(valid_files):
        original_filename = f.filename
        ext = original_filename.rsplit('.', 1)[1].lower()
        safe_temp = f'batch_{batch_id}_{idx:02d}.{ext}'
        filepath = os.path.join(session_folder, safe_temp)
        f.save(filepath)

        try:
            cf = color_filter_list if len(valid_files) == 1 else None
            extracted = _extract_raw(filepath, original_filename, glossary_ids, session_folder, color_filter=cf)
        except Exception as e:
            return jsonify({'error': f'Lỗi khi xử lý "{original_filename}": {str(e)}'}), 500

        batch_session_files.append({
            'idx': idx,
            'original_filename': original_filename,
            'display_name': original_filename,
            'filepath': filepath,
            'ext': ext,
        })
        file_extracted_list.append({
            'original_filename': original_filename,
            'extracted_data': extracted,
        })

    # ── Cross-file dedup ──────────────────────────────────────────────
    # value_to_refs: {value: {filename: [keys]}}
    value_to_refs: dict = {}
    for fe in file_extracted_list:
        fname = fe['original_filename']
        for key, value in fe['extracted_data'].items():
            if value not in value_to_refs:
                value_to_refs[value] = {}
            if fname not in value_to_refs[value]:
                value_to_refs[value][fname] = []
            value_to_refs[value][fname].append(key)

    # dedup_data: dedup_N → value (for translation)
    # cross_map:  dedup_N → {filename: [orig_keys]} (for injection)
    dedup_data: dict = {}
    cross_map: dict = {}
    for idx_d, (value, refs) in enumerate(value_to_refs.items(), 1):
        dk = f'dedup_{idx_d}'
        dedup_data[dk] = value
        cross_map[dk] = refs

    total_items = sum(len(fe['extracted_data']) for fe in file_extracted_list)
    unique_values = len(dedup_data)
    dedup_stats = {
        'total': total_items,
        'unique': unique_values,
        'saved': total_items - unique_values,
        'percent_saved': round((total_items - unique_values) * 100 / total_items) if total_items > 0 else 0,
    }

    # Chunk dedup_data into JSON parts
    CHUNK_SIZE = 300
    items_list = list(dedup_data.items())
    num_chunks = max(1, (unique_values + CHUNK_SIZE - 1) // CHUNK_SIZE)
    dedup_chunks = []
    for i in range(num_chunks):
        chunk = dict(items_list[i * CHUNK_SIZE:(i + 1) * CHUNK_SIZE])
        dedup_chunks.append({
            'name': f'batch_{batch_id}_dedup_part{i+1:02d}_of_{num_chunks:02d}.json',
            'content': json.dumps(chunk, ensure_ascii=False, indent=2),
        })

    # Build ZIP of dedup chunks for download
    zip_display_name = f'batch_{batch_id}_dedup.zip'
    zip_path = os.path.join(session_folder, f'batch_{batch_id}_extract.zip')
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_STORED) as zf:
        for c in dedup_chunks:
            zf.writestr(c['name'], c['content'])

    # Save cross_map to disk
    crossmap_path = os.path.join(session_folder, f'batch_{batch_id}_crossmap.json')
    with open(crossmap_path, 'w', encoding='utf-8') as f:
        json.dump(cross_map, f, ensure_ascii=False, indent=2)

    files_summary = [
        {'name': fe['original_filename'], 'items': len(fe['extracted_data'])}
        for fe in file_extracted_list
    ]

    session[f'batch_{batch_id}'] = {
        'batch_id': batch_id,
        'zip_path': zip_path,
        'zip_display_name': zip_display_name,
        'crossmap_path': crossmap_path,
        'files': batch_session_files,
    }

    return jsonify({
        'success': True,
        'batch_id': batch_id,
        'files': files_summary,
        'total_files': len(files_summary),
        'total_items': total_items,
        'zip_display_name': zip_display_name,
        'dedup_stats': dedup_stats,
        'dedup_chunks': dedup_chunks,
    })


@app.route('/download-batch-zip/<batch_id>', methods=['GET'])
@login_required
def download_batch_zip(batch_id):
    """Serve ZIP của batch extract."""
    key = f'batch_{batch_id}'
    if key not in session:
        return jsonify({'error': 'Batch session không tồn tại hoặc đã hết hạn.'}), 404
    info = session[key]
    zip_path = info.get('zip_path')
    if not zip_path or not os.path.exists(zip_path):
        return jsonify({'error': 'File ZIP không còn tồn tại.'}), 404
    response = send_file(zip_path, mimetype='application/zip')
    response = set_download_headers(response, info['zip_display_name'], 'batch_extract.zip')
    return response


@app.route('/batch-inject', methods=['POST'])
@login_required
def batch_inject():
    """
    Nạp bản dịch cho tất cả file trong batch, sử dụng cross-file dedup mapping.
    Input (form-data): batch_id + pasted_json_data (JSON string từ textarea)
    Output JSON: { success, files: [{token, display_name}] }
    """
    batch_id = request.form.get('batch_id', '').strip()
    key = f'batch_{batch_id}'
    if not batch_id or key not in session:
        return jsonify({'error': 'Batch session không tồn tại hoặc đã hết hạn. Vui lòng Extract lại.'}), 400

    batch_info = session[key]
    crossmap_path = batch_info.get('crossmap_path', '')
    if not crossmap_path or not os.path.exists(crossmap_path):
        return jsonify({'error': 'Cross-map không tồn tại. Vui lòng Extract lại.'}), 400

    # Nhận translated JSON từ paste
    translated_data: dict = {}
    pasted_raw = request.form.get('pasted_json_data', '')
    if pasted_raw:
        try:
            pasted_items = json.loads(pasted_raw)
            if isinstance(pasted_items, list):
                for item in pasted_items:
                    if isinstance(item, dict):
                        translated_data.update(item)
            elif isinstance(pasted_items, dict):
                translated_data.update(pasted_items)
        except json.JSONDecodeError as e:
            return jsonify({'error': f'JSON không hợp lệ: {str(e)}'}), 400
    # Also accept json_files[] for backward compat
    for jf in request.files.getlist('json_files'):
        if not jf.filename:
            continue
        raw = jf.stream.read()
        try:
            content = raw.decode('utf-8')
        except UnicodeDecodeError:
            content = raw.decode('utf-8-sig')
        try:
            translated_data.update(json.loads(content))
        except json.JSONDecodeError as e:
            return jsonify({'error': f'File "{jf.filename}" không phải JSON hợp lệ: {str(e)}'}), 400

    if not translated_data:
        return jsonify({'error': 'Không có dữ liệu JSON đã dịch.'}), 400

    # Load cross_map
    try:
        with open(crossmap_path, 'r', encoding='utf-8') as f:
            cross_map: dict = json.load(f)
    except Exception as e:
        return jsonify({'error': f'Lỗi đọc cross-map: {str(e)}'}), 500

    session_folder = get_session_folder()
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

    # Build per-file data from cross_map + translated_data
    per_file_data: dict = {info['original_filename']: {} for info in batch_info['files']}
    for dk, refs in cross_map.items():
        if dk not in translated_data:
            continue
        translated_value = translated_data[dk]
        for fname, orig_keys in refs.items():
            if fname in per_file_data:
                for ok in orig_keys:
                    per_file_data[fname][ok] = translated_value

    error_details: list = []
    result_files: list = []

    for source_info in batch_info['files']:
        source_name = source_info['original_filename']
        source_filepath = source_info['filepath']
        ext = source_info['ext']
        file_json_data = per_file_data.get(source_name, {})

        if not file_json_data:
            error_details.append(f'"{source_name}": không có dữ liệu dịch')
            continue
        if not os.path.exists(source_filepath):
            error_details.append(f'"{source_name}": file gốc đã hết hạn')
            continue

        out_display = f'{os.path.splitext(source_name)[0]}_translated.{ext}'
        out_path = os.path.join(session_folder, f'bout_{batch_id}_{timestamp}_{source_info["idx"]:02d}.{ext}')

        try:
            if ext == 'xlsx':
                inject_xlsx_shapes(source_filepath, out_path, file_json_data)
            elif ext == 'pptx':
                prs = inject_text_to_pptx(source_filepath, file_json_data)
                prs.save(out_path)
                del prs
            elif ext == 'docx':
                doc = inject_text_to_docx(source_filepath, file_json_data)
                doc.save(out_path)
                del doc
            else:
                error_details.append(f'"{source_name}": định dạng .{ext} không hỗ trợ')
                continue

            if os.path.exists(out_path):
                token = uuid.uuid4().hex[:14]
                session[f'injected_{token}'] = {
                    'path': out_path,
                    'display_name': out_display,
                }
                result_files.append({'token': token, 'display_name': out_display})
            else:
                error_details.append(f'"{source_name}": không tạo được file output')
        except Exception as e:
            error_details.append(f'"{source_name}": {str(e)}')

    if not result_files:
        return jsonify({'error': 'Không inject được file nào. ' + '; '.join(error_details)}), 500

    return jsonify({
        'success': True,
        'files': result_files,
        'errors': error_details,
    })


@app.route('/download-injected/<token>', methods=['GET'])
@login_required
def download_injected(token):
    """Tải xuống file đã inject theo token."""
    token_key = f'injected_{token}'
    if token_key not in session:
        return jsonify({'error': 'File không tồn tại hoặc đã hết hạn.'}), 404
    info = session[token_key]
    file_path = info.get('path', '')
    display_name = info.get('display_name', 'translated_file')
    if not file_path or not os.path.exists(file_path):
        session.pop(token_key, None)
        return jsonify({'error': 'File không còn tồn tại.'}), 404

    ext = file_path.rsplit('.', 1)[-1].lower()
    mime_map = {
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    }
    response = send_file(file_path, mimetype=mime_map.get(ext, 'application/octet-stream'))
    response = set_download_headers(response, display_name, f'translated.{ext}')

    session.pop(token_key, None)

    @response.call_on_close
    def _cleanup():
        import time as _t, gc
        gc.collect()
        _t.sleep(0.1)
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass

    return response


@app.route('/batch-inject-one', methods=['POST'])
@login_required
def batch_inject_one():
    """
    Nạp bản dịch cho một file trong batch, sử dụng cross-file dedup mapping.
    Input (form-data): batch_id + source_filename + json_files[]
    Output: file đã inject bản dịch (trực tiếp)
    """
    batch_id = request.form.get('batch_id', '').strip()
    source_filename = request.form.get('source_filename', '').strip()
    key = f'batch_{batch_id}'

    if not batch_id or key not in session:
        return jsonify({'error': 'Batch session không tồn tại hoặc đã hết hạn.'}), 400
    if not source_filename:
        return jsonify({'error': 'Thiếu tên file gốc (source_filename).'}), 400

    batch_info = session[key]
    session_files = {info['original_filename']: info for info in batch_info['files']}

    if source_filename not in session_files:
        return jsonify({'error': f'File "{source_filename}" không có trong batch session.'}), 400

    source_info = session_files[source_filename]
    source_filepath = source_info['filepath']
    if not os.path.exists(source_filepath):
        return jsonify({'error': f'File gốc "{source_filename}" không còn tồn tại.'}), 400

    json_files_upload = request.files.getlist('json_files')
    if not any(f.filename for f in json_files_upload):
        return jsonify({'error': 'Cần upload ít nhất 1 file JSON.'}), 400

    translated_data: dict = {}
    for jf in json_files_upload:
        if not jf.filename:
            continue
        raw = jf.stream.read()
        try:
            content = raw.decode('utf-8')
        except UnicodeDecodeError:
            content = raw.decode('utf-8-sig')
        try:
            translated_data.update(json.loads(content))
        except json.JSONDecodeError as e:
            return jsonify({'error': f'File "{jf.filename}" không phải JSON hợp lệ: {str(e)}'}), 400

    # Use cross_map to reconstruct per-file data for this source file
    crossmap_path = batch_info.get('crossmap_path', '')
    json_data: dict = {}
    if crossmap_path and os.path.exists(crossmap_path):
        try:
            with open(crossmap_path, 'r', encoding='utf-8') as f:
                cross_map = json.load(f)
            for dk, refs in cross_map.items():
                if source_filename in refs and dk in translated_data:
                    for orig_key in refs[source_filename]:
                        json_data[orig_key] = translated_data[dk]
        except Exception:
            json_data = translated_data  # fallback: use translated directly
    else:
        json_data = translated_data  # no cross_map → use translated directly

    if not json_data:
        return jsonify({'error': f'Không tìm thấy dữ liệu dịch cho file "{source_filename}".'}), 400

    session_folder = get_session_folder()
    ext = source_info['ext']
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    out_display = f'{os.path.splitext(source_filename)[0]}_translated.{ext}'
    out_path = os.path.join(session_folder, f'batchone_{batch_id}_{timestamp}.{ext}')

    mime_map = {
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    }

    try:
        if ext == 'xlsx':
            inject_xlsx_shapes(source_filepath, out_path, json_data)
        elif ext == 'pptx':
            prs = inject_text_to_pptx(source_filepath, json_data)
            prs.save(out_path)
            del prs
        elif ext == 'docx':
            doc = inject_text_to_docx(source_filepath, json_data)
            doc.save(out_path)
            del doc
        else:
            return jsonify({'error': f'Định dạng .{ext} không hỗ trợ.'}), 400
    except Exception as e:
        return jsonify({'error': f'Lỗi khi inject: {str(e)}'}), 500

    response = send_file(out_path, mimetype=mime_map.get(ext, 'application/octet-stream'))
    response = set_download_headers(response, out_display, f'download.{ext}')

    @response.call_on_close
    def _cleanup():
        import time as _t, gc
        gc.collect()
        _t.sleep(0.1)
        try:
            if os.path.exists(out_path):
                os.remove(out_path)
        except Exception:
            pass

    return response


if __name__ == '__main__':
    # Chạy ứng dụng Flask ở chế độ debug
    #app.run(host='0.0.0.0', port=5000)
    app.run(debug=True,host='0.0.0.0', port=5017)
